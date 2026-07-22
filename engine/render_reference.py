"""
Faithful PPTX-slide -> HTML renderer for the reference library.

Reference slides are REUSED VERBATIM: word-for-word, image-for-image, style-for-style.
The ONLY sanctioned transform is colour conversion to the brand guide (dark<->light
theme). So instead of re-authoring a slide into plan blocks and re-styling it (which
loses fidelity), we reconstruct the slide as absolute-positioned HTML on the 1280x720
canvas: exact text runs (font/size/weight/italic/colour/align resolved through the PPTX
inheritance chain), images at their exact box with PowerPoint crop/stretch semantics,
solid/gradient fills, connector lines, and — for shapes that cannot be faithfully rebuilt
(SmartArt, freeforms, connectors-with-arrowheads, charts, groups) — a crisp COM-exported
raster embedded at the shape box. A verbatim-locked slide never silently omits a visual.

    from render_reference import render_slide, load_theme_map, load_font_scheme, rasterize_shapes
    raster = rasterize_shapes(pptx_path, idx0)            # {shape_idx0: data_uri}  (needs COM)
    html = render_slide(prs, idx0, theme_map, font_scheme, theme='dark', raster_shapes=raster)

Produced HTML is a `.stage` inner fragment; build.py drops it in unchanged for a reused
slide (locked — REFINE/SWEEP must not content-edit it).

Fixes tracked: RC1 all-or-nothing recolour · RC2 crop/stretch image semantics ·
RC3 group/diagram recursion + raster fallback · RC4 theme fonts + inherited bold/size + a:br.
"""
import base64, html as _html, os, re, zipfile
from xml.etree import ElementTree as ET

EMU_PER_PT = 12700
CANVAS_W, CANVAS_H = 1280, 720
A = "http://schemas.openxmlformats.org/drawingml/2006/main"

LOGO_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "assets", "logos")


# ============================================================ brand-logo SVG inlining
# Shared by the reference-review builder AND ingest (_render_verbatim) so catalog renders
# match the exec-approved review deck exactly. Handles: XML-prolog strip, per-file class
# de-collision, id namespacing for gradients/clippaths (id collisions made four logos
# invisible — cycle-9), and a DARK variant that lifts neutral GREY letterforms (flat fills,
# GRADIENT stops, AND strokes) so they read on navy (W1d).
def _grey_lift(hexv, thresh_sat=28, max_lum=0.60, to="e6ecf3"):
    """Return the light replacement hex if hexv is a neutral, dark-ish grey, else None.
    Colored (blue/orange/navy) values are left untouched — mx-mn guards saturation."""
    try:
        r, g, b = int(hexv[0:2], 16), int(hexv[2:4], 16), int(hexv[4:6], 16)
    except Exception:
        return None
    mx, mn = max(r, g, b), min(r, g, b)
    lum = (0.2126 * r + 0.7152 * g + 0.0722 * b) / 255.0
    if (mx - mn) <= thresh_sat and lum < max_lum:
        return to
    return None


def logo_svg(name, dark=False):
    """Inline an SVG brand logo from assets/logos/ for embedding in HTML.
    dark=True lifts grey letterforms so they survive on navy backgrounds."""
    path = name if os.path.isabs(name) else os.path.join(LOGO_DIR, name)
    t = open(path, encoding="utf-8").read()
    t = re.sub(r'<\?xml[^>]*\?>', "", t).strip()             # XML prolog breaks HTML inlining
    slug = re.sub(r'[^a-z0-9]+', '', os.path.basename(name).lower())
    t = t.replace("cls-", "c%s-" % slug)                      # de-collide per-file class names
    for _id in set(re.findall(r'id="([^"]+)"', t)):          # namespace ALL ids + refs
        t = t.replace('id="%s"' % _id, 'id="%s-%s"' % (_id, slug))
        t = t.replace('url(#%s)' % _id, 'url(#%s-%s)' % (_id, slug))
        t = t.replace('href="#%s"' % _id, 'href="#%s-%s"' % (_id, slug))
    if dark:
        def _f(m):
            rep = _grey_lift(m.group(1))
            return m.group(0).replace(m.group(1), rep) if rep else m.group(0)
        t = re.sub(r'fill:\s*#([0-9a-fA-F]{6})', _f, t)      # flat fills (style block)
        t = re.sub(r'fill="#([0-9a-fA-F]{6})"', _f, t)       # flat fills (attribute)
        t = re.sub(r'stop-color[:=]\s*"?#([0-9a-fA-F]{6})"?', _f, t)   # W1d: gradient stops
        t = re.sub(r'stroke:\s*#([0-9a-fA-F]{6})', _f, t)    # W1d: strokes
        t = re.sub(r'stroke="#([0-9a-fA-F]{6})"', _f, t)
    m = re.search(r'<svg[^>]*>', t)
    tag = m.group(0)
    tag2 = re.sub(r'\swidth="[^"]*"', "", tag)
    tag2 = re.sub(r'\sheight="[^"]*"', "", tag2)
    tag2 = tag2.replace("<svg", '<svg width="100%" height="100%" preserveAspectRatio="xMidYMid meet"', 1)
    return t.replace(tag, tag2, 1)


REF_ASSET_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                             "layouts", "reference-library", "assets")


def img_swap(name):
    """Inline a curated PNG/JPEG asset (reference-library/assets/) as an <img> that fills
    its shape box — a theme-independent RASTER counterpart to logo_svg(). Used by
    curation `image_swaps` to replace a source picture whose baked pixels are wrong (e.g.
    the s08 glass cloud carried a flat ground-shadow that read badly on dark — cycle-10
    pin-5). object-fit:contain keeps the cut-out at the same position/scale in the box."""
    path = name if os.path.isabs(name) else os.path.join(REF_ASSET_DIR, name)
    with open(path, "rb") as fh:
        blob = fh.read()
    ext = os.path.splitext(path)[1].lower()
    ctype = {".jpg": "image/jpeg", ".jpeg": "image/jpeg",
             ".svg": "image/svg+xml"}.get(ext, "image/png")
    uri = "data:%s;base64,%s" % (ctype, base64.b64encode(blob).decode("ascii"))
    return ('<img src="%s" alt="" style="width:100%%;height:100%%;object-fit:contain;'
            'display:block">' % uri)


# ============================================================ theme colour + font
def load_theme_map(pptx_path):
    """Read ppt/theme/theme1.xml clrScheme -> {'dk1':'000000','lt1':'FFFFFF',...}."""
    out = {}
    try:
        with zipfile.ZipFile(pptx_path) as z:
            name = next((n for n in z.namelist() if re.match(r"ppt/theme/theme\d+\.xml$", n)), None)
            if not name:
                return out
            root = ET.fromstring(z.read(name))
        scheme = root.find(".//{%s}clrScheme" % A)
        if scheme is None:
            return out
        for child in scheme:
            key = child.tag.split("}")[1]
            srgb = child.find("{%s}srgbClr" % A)
            sysc = child.find("{%s}sysClr" % A)
            if srgb is not None:
                out[key] = srgb.get("val")
            elif sysc is not None:
                out[key] = sysc.get("lastClr") or sysc.get("val")
    except Exception:
        pass
    return out


def load_font_scheme(pptx_path):
    """major/minor latin typefaces from theme1.xml fontScheme (RC4)."""
    fs = {"major": "Montserrat", "minor": "Montserrat"}
    try:
        with zipfile.ZipFile(pptx_path) as z:
            name = next((n for n in z.namelist() if re.match(r"ppt/theme/theme\d+\.xml$", n)), None)
            root = ET.fromstring(z.read(name))
        sc = root.find(".//{%s}fontScheme" % A)
        for role, key in (("major", "majorFont"), ("minor", "minorFont")):
            lat = sc.find("{%s}%s/{%s}latin" % (A, key, A))
            if lat is not None and lat.get("typeface"):
                fs[role] = lat.get("typeface")
    except Exception:
        pass
    return fs


_THEME_KEY = {
    "DARK_1": "dk1", "LIGHT_1": "lt1", "DARK_2": "dk2", "LIGHT_2": "lt2",
    "ACCENT_1": "accent1", "ACCENT_2": "accent2", "ACCENT_3": "accent3",
    "ACCENT_4": "accent4", "ACCENT_5": "accent5", "ACCENT_6": "accent6",
    "HYPERLINK": "hlink", "FOLLOWED_HYPERLINK": "folHlink",
    "TEXT_1": "dk1", "TEXT_2": "dk2", "BACKGROUND_1": "lt1", "BACKGROUND_2": "lt2",
}
# scheme keys as they appear inside <a:schemeClr val="..."> (lowercase, tx1/bg1 aliases)
_SCHEME_ALIAS = {"tx1": "dk1", "tx2": "dk2", "bg1": "lt1", "bg2": "lt2"}


def _scheme_hex(val, theme_map):
    if not val:
        return None
    key = _SCHEME_ALIAS.get(val, val)
    return theme_map.get(key) or theme_map.get(_THEME_KEY.get(val.upper(), ""))


def _apply_mods(hexv, clr_el):
    """Apply DrawingML colour modifiers (lumMod/lumOff/tint/shade/satMod) to a hex colour.
    These were IGNORED, so light-tinted brand fills rendered fully saturated (loop cycle-1:
    s10's soft light-blue panel came out bright cyan)."""
    try:
        import colorsys
        r, g, b = (v / 255.0 for v in _rgb(hexv))
        def pct(tag):
            e = clr_el.find("{%s}%s" % (A, tag))
            return (int(e.get("val")) / 100000.0) if (e is not None and e.get("val")) else None
        tint, shade = pct("tint"), pct("shade")
        lumMod, lumOff, satMod = pct("lumMod"), pct("lumOff"), pct("satMod")
        if tint is not None:                       # toward white
            r, g, b = (c * tint + (1 - tint) for c in (r, g, b))
        if shade is not None:                      # toward black
            r, g, b = (c * shade for c in (r, g, b))
        if lumMod is not None or lumOff is not None or satMod is not None:
            h, l, s = colorsys.rgb_to_hls(r, g, b)
            if satMod is not None:
                s = min(1.0, s * satMod)
            if lumMod is not None:
                l = l * lumMod
            if lumOff is not None:
                l = min(1.0, max(0.0, l + lumOff))
            r, g, b = colorsys.hls_to_rgb(h, l, s)
        return "#%02x%02x%02x" % (int(r * 255 + .5), int(g * 255 + .5), int(b * 255 + .5))
    except Exception:
        return hexv


def _solidfill_color(sf, theme_map):
    """(css_color, hex_or_None) from an <a:solidFill>, applying scheme lookup, modifiers
    and alpha. The one true fill-colour resolver — python-pptx fore_color ignores mods."""
    if sf is None:
        return None, None
    srgb = sf.find("{%s}srgbClr" % A); sch = sf.find("{%s}schemeClr" % A)
    clr = srgb if srgb is not None else sch
    if clr is None:
        return None, None
    hexc = ("#" + srgb.get("val")) if srgb is not None else (
        ("#" + _scheme_hex(sch.get("val", ""), theme_map)) if _scheme_hex(sch.get("val", ""), theme_map) else None)
    if not hexc:
        return None, None
    hexc = _apply_mods(hexc, clr)
    a = clr.find("{%s}alpha" % A)
    if a is not None and a.get("val"):
        alpha = int(a.get("val")) / 100000.0
        if alpha < 0.995:
            r, g, b = _rgb(hexc)
            return "rgba(%d,%d,%d,%.3f)" % (r, g, b, alpha), hexc
    return hexc, hexc


def _sx(prs):
    return CANVAS_W / prs.slide_width


def _sy(prs):
    return CANVAS_H / prs.slide_height


_ALIGN = {"CENTER": "center", "RIGHT": "right", "LEFT": "left", "JUSTIFY": "justify"}


# ============================================================ brand tokens + recolour (RC1)
DARK_BG = "#0B1E3A"        # brand navy — generated dark background
LIGHT_BG = "#FFFFFF"
DARK_PANEL = "#0d1829"     # panel surface on dark (near-white fills map here)
LIGHT_PANEL = "#f5f8fc"    # panel surface on light
LIGHT_TEXT = "#FFFFFF"
DARK_INK = "#0F1419"
# brand accents that PASS THROUGH recolour unchanged
_BRAND = ["00B2E3", "4FC3F7", "00B2E2", "E17126", "76A3B2", "ED1651", "0B2A4A", "0B1E3A", "060C1A"]


def _rgb(hexv):
    h = (hexv or "").lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _luminance(hexv):
    try:
        r, g, b = _rgb(hexv)
        return (0.2126 * r + 0.7152 * g + 0.0722 * b) / 255.0
    except Exception:
        return 0.5


def _is_brandish(hexv):
    try:
        r, g, b = _rgb(hexv)
    except Exception:
        return False
    for br in _BRAND:
        R, G, B = _rgb(br)
        if abs(r - R) + abs(g - G) + abs(b - B) <= 60:   # near a brand accent
            return True
    return False


def _map_color(hexv, target):
    """Surface mapping (cycle-7 rethink, per Rowan): DARK theme surfaces live in the NAVY
    family only — no green/teal/odd-hue panels. Brand sky/copper pass through for accents.
    Tinted light fills all land on the navy panel ramp; sky-family tints may keep a deep-sky
    cast. Light theme unchanged-ish."""
    if not hexv or hexv.startswith("rgba"):
        return hexv
    if _is_brandish(hexv):
        return hexv
    try:
        import colorsys
        r, g, b = (v / 255.0 for v in _rgb(hexv))
        h, l, sat = colorsys.rgb_to_hls(r, g, b)
    except Exception:
        return hexv
    def out(h2, l2, s2):
        r2, g2, b2 = colorsys.hls_to_rgb(h2, max(0.0, min(1.0, l2)), max(0.0, min(1.0, s2)))
        return "#%02x%02x%02x" % (int(r2 * 255 + .5), int(g2 * 255 + .5), int(b2 * 255 + .5))
    NAVY_H, SKY_H = 0.583, 0.55
    if target == "dark":
        sky_family = 0.50 <= h <= 0.62
        if sat < 0.12:                                    # neutral surfaces -> navy ramp
            if l >= 0.85: return out(NAVY_H, 0.115, 0.55)
            if l >= 0.60: return out(NAVY_H, 0.16, 0.45)
            if l <= 0.30: return "#F2F7FC"
            return out(NAVY_H, 0.55, 0.15)
        if l >= 0.55:
            # cycle-10 pins 7/12/13/15: a tinted light fill keeps its HUE as a dark
            # accent surface (number-circle halos, card gradients, the s10 panel) —
            # flattening every tint to navy washed the dark theme out. Sky-family
            # tints land on a clearly-blue elevated panel ('lighter blue, higher
            # contrast'); other hues drop to a dark tone of themselves.
            if sky_family:
                return out(SKY_H, 0.30, 0.52)
            return out(h, 0.33, 0.58)                     # cycle-10 pin-15: bolder tinted
            #                                               surfaces (s18 card gradients) —
            #                                               0.26/0.45 read washed-out on navy
        if sat >= 0.30:                                   # vivid accents read on navy — keep
            return out(h, max(l, 0.32), sat)
        if sky_family:                                    # deep sky tones keep their cast
            return out(h, max(0.30, l), sat)
        return out(NAVY_H, 0.16, 0.45)                    # muted odd-hue surface -> navy
    else:
        if sat < 0.12:
            if l <= 0.30: return "#13202f"
            if l >= 0.85: return "#f4f7fb"
            return out(0.0, 0.35, 0.0)
        if l <= 0.45: return out(h, 0.88, max(0.25, sat * 0.6))
        return out(h, min(0.92, l), sat)


def _map_text_color(col, target):
    """TEXT mapping (cycle-7 rethink, per Rowan): on dark, WHITE is the primary text
    colour; only true brand accents (sky, copper) stay coloured. No hue-walking — a green
    or random-hue run becomes white, not a lighter green."""
    if target == "dark":
        col = col or "#F2F7FC"
        if _is_brandish(col) and _contrast_ratio(col, DARK_BG) >= 3.0:
            return col                                    # sky / copper accent survives
        if _contrast_ratio(col, DARK_BG) >= 4.5 and _luminance(col) > 0.62:
            return col                                    # already light -> keep
        return "#F2F7FC"                                  # everything else: white primary
    col = col or "#13202f"
    if _contrast_ratio(col, LIGHT_BG) >= 4.5:
        return col
    if _is_brandish(col):
        return "#0b6ea8"
    return "#13202f"


def _contrast_ratio(fg, bg):
    def _lin(c):
        c /= 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    def _rl(hexv):
        r, g, b = _rgb(hexv)
        return 0.2126 * _lin(r) + 0.7152 * _lin(g) + 0.0722 * _lin(b)
    try:
        l1, l2 = _rl(fg), _rl(bg)
        hi, lo = max(l1, l2), min(l1, l2)
        return (hi + 0.05) / (lo + 0.05)
    except Exception:
        return 21.0


# ============================================================ background + images
def slide_background(slide, theme_map):
    chain = [slide]
    try:
        chain.append(slide.slide_layout)
        chain.append(slide.slide_layout.slide_master)
    except Exception:
        pass
    for src in chain:
        try:
            f = src.background.fill
        except Exception:
            continue
        if f.type == 1:
            fc = f.fore_color
            try:
                if fc.type == 1:
                    return "#" + str(fc.rgb)
            except Exception:
                pass
            key = _THEME_KEY.get(getattr(fc.theme_color, "name", "") or "")
            if key and theme_map.get(key):
                return "#" + theme_map[key]
    return "#FFFFFF"


def native_theme_of(prs, slide, theme_map):
    try:
        import io
        from PIL import Image, ImageStat
        for sh in slide.shapes:
            if (sh.shape_type == 13 and sh.width and sh.height and sh.left is not None
                    and sh.width >= prs.slide_width * 0.9 and sh.height >= prs.slide_height * 0.9):
                im = Image.open(io.BytesIO(sh.image.blob)).convert("L")
                return "dark" if (ImageStat.Stat(im).mean[0] / 255.0) < 0.5 else "light"
    except Exception:
        pass
    return "light" if _luminance(slide_background(slide, theme_map)) >= 0.5 else "dark"


def _process_image(blob, content_type, target):
    """Brighten (light variant) or darken (dark variant) a full-bleed background image."""
    try:
        import io
        from PIL import Image, ImageEnhance
        im = Image.open(io.BytesIO(blob)).convert("RGB")
        if target == "light":
            im = ImageEnhance.Brightness(im).enhance(1.35)
            im = ImageEnhance.Color(im).enhance(0.8)
            im = Image.blend(im, Image.new("RGB", im.size, (255, 255, 255)), 0.16)
        else:
            im = ImageEnhance.Brightness(im).enhance(0.45)
            im = Image.blend(im, Image.new("RGB", im.size, (8, 20, 40)), 0.35)
        buf = io.BytesIO(); im.save(buf, format="JPEG", quality=85)
        return "data:image/jpeg;base64," + base64.b64encode(buf.getvalue()).decode()
    except Exception:
        return None


def _embed_image(blob, content_type):
    """Embed an image as a data URI. RC5: NEVER JPEG-flatten an image that has an alpha
    channel (a transparent brand mark would get a white box); keep it PNG. Only large
    PHOTOS (opaque, or huge) are recompressed to JPEG — exact pixels don't matter there."""
    try:
        import io
        from PIL import Image
        im = Image.open(io.BytesIO(blob))
        has_alpha = ("A" in im.getbands()) or (im.mode == "P" and "transparency" in im.info)
        if has_alpha:
            im = im.convert("RGBA")
            buf = io.BytesIO(); im.save(buf, format="PNG", optimize=True)
            return "data:image/png;base64," + base64.b64encode(buf.getvalue()).decode()
        if len(blob) > 300_000:            # opaque photo -> JPEG to keep the file sane
            im = im.convert("RGB")
            buf = io.BytesIO(); im.save(buf, format="JPEG", quality=84)
            return "data:image/jpeg;base64," + base64.b64encode(buf.getvalue()).decode()
    except Exception:
        pass
    return "data:%s;base64,%s" % (content_type, base64.b64encode(blob).decode("ascii"))


def _src_rect(sh):
    """PowerPoint <a:srcRect> crop (l/t/r/b in 1000ths of a percent) -> CSS background-size
    + background-position that reproduces the exact crop. Returns (size, pos) strings.
    No srcRect -> full stretch (background-size:100% 100%), which is PowerPoint's DEFAULT
    picture fill semantic — NOT object-fit:contain (that letterboxed images: H1/H4/H6)."""
    try:
        # blipFill on a <p:pic> is in the PRESENTATIONML namespace; on a picture-filled
        # autoshape it's drawingml. The a:-only find NEVER matched pics, so srcRect crops
        # were silently skipped — the root cause of every 'stretched image' review pin.
        P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        blip = sh._element.find("{%s}blipFill" % P)
        if blip is None:
            blip = sh._element.find(".//{%s}blipFill" % A)
        sr = blip.find("{%s}srcRect" % A) if blip is not None else None
    except Exception:
        sr = None
    if sr is None or not any(sr.get(k) for k in ("l", "t", "r", "b")):
        return None  # no crop
    l = int(sr.get("l", "0")) / 100000.0
    t = int(sr.get("t", "0")) / 100000.0
    r = int(sr.get("r", "0")) / 100000.0
    b = int(sr.get("b", "0")) / 100000.0
    return (l, t, r, b)


def _blip_fx(shape, blob, theme_map, recolor=None):
    """Apply blip-level picture effects PowerPoint bakes at render time (cycle-5):
    <a:clrChange> ('set transparent colour' — a white-bg logo hid the chart axis behind
    it) and <a:duotone> (grayscale -> two-colour ramp, s18 icons). Returns PNG bytes or
    None if no effect applies."""
    try:
        P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        blipFill = shape._element.find("{%s}blipFill" % P)
        if blipFill is None:
            blipFill = shape._element.find(".//{%s}blipFill" % A)
        blip = blipFill.find("{%s}blip" % A) if blipFill is not None else None
        if blip is None:
            return None
        cc = blip.find("{%s}clrChange" % A)
        duo = blip.find("{%s}duotone" % A)
        if cc is None and duo is None:
            return None
        import io
        from PIL import Image
        im = Image.open(io.BytesIO(blob)).convert("RGBA")
        if cc is not None:
            frm = cc.find(".//{%s}clrFrom/{%s}srgbClr" % (A, A))
            to = cc.find(".//{%s}clrTo/{%s}srgbClr" % (A, A))
            to_alpha = to.find("{%s}alpha" % A) if to is not None else None
            if frm is not None and to_alpha is not None and int(to_alpha.get("val", "100000")) == 0:
                fr, fg, fb = _rgb("#" + frm.get("val"))
                px = im.load(); w, h = im.size
                from collections import deque
                TOL = 28
                def _match(xx, yy):
                    r, g, b, a = px[xx, yy]
                    return a > 0 and abs(r - fr) <= TOL and abs(g - fg) <= TOL and abs(b - fb) <= TOL
                seen = bytearray(w * h); q = deque()
                for xx in range(w):
                    for yy in (0, h - 1):
                        if _match(xx, yy) and not seen[yy * w + xx]:
                            seen[yy * w + xx] = 1; q.append((xx, yy))
                for yy in range(h):
                    for xx in (0, w - 1):
                        if _match(xx, yy) and not seen[yy * w + xx]:
                            seen[yy * w + xx] = 1; q.append((xx, yy))
                while q:
                    xx, yy = q.popleft()
                    r, g, b, a = px[xx, yy]
                    px[xx, yy] = (r, g, b, 0)
                    for nx, ny in ((xx-1, yy), (xx+1, yy), (xx, yy-1), (xx, yy+1)):
                        if 0 <= nx < w and 0 <= ny < h and not seen[ny * w + nx] and _match(nx, ny):
                            seen[ny * w + nx] = 1; q.append((nx, ny))
        if duo is not None:
            cols = []
            for ch in duo:
                tag = ch.tag.split("}")[1]
                if tag == "srgbClr":
                    cols.append(_rgb("#" + ch.get("val")))
                elif tag == "schemeClr":
                    hx = _scheme_hex(ch.get("val", ""), theme_map)
                    if hx:
                        cols.append(_rgb(_apply_mods("#" + hx, ch)))
                elif tag == "prstClr":
                    cols.append((0, 0, 0) if ch.get("val") == "black" else (255, 255, 255))
            if len(cols) == 2:
                if recolor == "dark":                     # cycle-9 pin-7: MAX contrast line art
                    c0_lum = _luminance("#%02x%02x%02x" % cols[0])
                    c1_lum = _luminance("#%02x%02x%02x" % cols[1])
                    if c0_lum <= c1_lum:
                        cols = [(255, 255, 255), (11, 30, 58)]
                    else:
                        cols = [(11, 30, 58), (255, 255, 255)]
                (r1, g1, b1), (r2, g2, b2) = cols
                gray = im.convert("L")
                px, gpx = im.load(), gray.load()
                w, h = im.size
                for yy in range(h):
                    for xx in range(w):
                        t = gpx[xx, yy] / 255.0
                        a = px[xx, yy][3]
                        px[xx, yy] = (int(r1 + (r2 - r1) * t), int(g1 + (g2 - g1) * t),
                                      int(b1 + (b2 - b1) * t), a)
        import io as _io
        buf = _io.BytesIO(); im.save(buf, format="PNG", optimize=True)
        return buf.getvalue()
    except Exception:
        return None


def _pic_html(uri, x, y, ww, hh, crop):
    """Picture at its frame with PowerPoint crop semantics, incl. NEGATIVE srcRect values
    (window larger than the source image). Wrapper clips; the inner img is scaled and
    offset in percentages of the frame — exact for any l/t/r/b sign."""
    if not crop:
        return ('<div style="%soverflow:hidden;">'
                '<img src="%s" alt="" style="position:absolute;left:0;top:0;width:100%%;height:100%%;display:block"></div>'
                % (_pos_style(x, y, ww, hh), uri))
    l, t, r, b = crop
    vw = max(1e-6, 1 - l - r)   # visible source window mapped onto the frame
    vh = max(1e-6, 1 - t - b)
    w_pct = 100.0 / vw
    h_pct = 100.0 / vh
    left = -(l / vw) * 100.0
    top = -(t / vh) * 100.0
    return ('<div style="%soverflow:hidden;">'
            '<img src="%s" alt="" style="position:absolute;left:%.2f%%;top:%.2f%%;width:%.2f%%;height:%.2f%%;display:block">'
            '</div>' % (_pos_style(x, y, ww, hh), uri, left, top, w_pct, h_pct))


# ============================================================ COM shape rasterization (RC3)
def unrenderable_pics(prs, index0):
    """Indices of pictures python-pptx cannot expose a blob for (linked/graphic parts) —
    these must be COM-rasterized or they'd vanish (review pin-12)."""
    out = set()
    for i, sh in enumerate(prs.slides[index0].shapes):
        if sh.shape_type == 13:
            try:
                sh.image.blob
            except Exception:
                out.add(i)
    return out


def rasterize_shapes(pptx_path, index0, only_idx=None, scale=3, extra_idx=None):
    """Rasterize complex top-level shapes (groups, charts, SmartArt, freeforms) of slide
    index0 to PNG data URIs. Returns {shape_idx0: data_uri}. Best-effort: {} if COM is
    unavailable (the renderer then recurses/omits with a warning).

    Method (ISOLATED per shape — review fix for the double-render/'33' defect): work on a
    TEMP COPY of the deck; for each target shape hide every OTHER shape, export the slide,
    crop the target's box, restore visibility. The crop keeps the composed on-slide
    placement (reliable for rotated groups/radial labels) but can no longer capture
    NEIGHBOURING content — previously an overlapping neighbour's text was baked into the
    crop AND drawn natively, printing twice ('3' -> '33', garbled captions). scale=3 for
    crisp text inside rasters."""
    out = {}
    try:
        import win32com.client, pythoncom, os, tempfile, io, shutil
        from pptx import Presentation
        from PIL import Image
    except Exception:
        return out
    try:
        prs = Presentation(pptx_path)
        slide = prs.slides[index0]
        sw, sh_emu = prs.slide_width, prs.slide_height
        targets = {}
        for i, shp in enumerate(slide.shapes):
            st = shp.shape_type
            if st == 6 and is_simple_group(shp):
                continue                     # native recursion (cycle-3 pin-2: real text, no raster)
            is_complex = (st == 6) or (st == 3)
            try:
                is_complex = is_complex or (st is not None and int(st) in (7, 24, 21))
            except Exception:
                pass
            if only_idx is not None:
                is_complex = i in only_idx
            if extra_idx and i in extra_idx:
                is_complex = True
            if is_complex and shp.left is not None and shp.width and shp.height:
                targets[i] = (shp.left, shp.top, shp.width, shp.height)
        if not targets:
            return out
    except Exception:
        return out
    pythoncom.CoInitialize()
    app = None
    tmp = None
    try:
        tmp = tempfile.mkdtemp()
        work = os.path.join(tmp, "work.pptx")
        shutil.copyfile(os.path.abspath(pptx_path), work)   # never touch the source
        app = win32com.client.Dispatch("PowerPoint.Application")
        pres = app.Presentations.Open(work, ReadOnly=False, Untitled=False, WithWindow=False)
        com_slide = pres.Slides(index0 + 1)
        n_shapes = com_slide.Shapes.Count
        for i, (l, t, w, h) in targets.items():
            # isolate: only the target visible (COM Shapes is 1-based, document order —
            # same order python-pptx iterates)
            for j in range(1, n_shapes + 1):
                com_slide.Shapes(j).Visible = (j == i + 1)
            png = os.path.abspath(os.path.join(tmp, "shape%d.png" % i))
            com_slide.Export(png, "PNG", CANVAS_W * scale, CANVAS_H * scale)
            page = Image.open(png).convert("RGBA")
            pw, ph = page.size
            x0 = int(l / sw * pw); y0 = int(t / sh_emu * ph)
            x1 = int((l + w) / sw * pw); y1 = int((t + h) / sh_emu * ph)
            crop = page.crop((max(0, x0), max(0, y0), min(pw, x1), min(ph, y1)))
            buf = io.BytesIO(); crop.save(buf, format="PNG")
            out[i] = "data:image/png;base64," + base64.b64encode(buf.getvalue()).decode("ascii")
        for j in range(1, n_shapes + 1):                     # restore (copy is discarded anyway)
            com_slide.Shapes(j).Visible = True
        pres.Close()
    except Exception:
        pass
    finally:
        try:
            if app is not None:
                app.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()
        try:
            if tmp:
                import shutil as _sh
                _sh.rmtree(tmp, ignore_errors=True)
        except Exception:
            pass
    return out



def _autofit_scale(shape):
    """<a:normAutofit fontScale lnSpcReduction> — PowerPoint SHRINKS text to fit its box;
    ignoring it rendered long titles oversized (cycle-3 pins: header too big / wrong wrap)."""
    try:
        bp = shape.text_frame._txBody.find("{%s}bodyPr" % A)
        na = bp.find("{%s}normAutofit" % A) if bp is not None else None
        if na is None:
            return 1.0, 0.0
        fs = int(na.get("fontScale", "100000")) / 100000.0
        lr = int(na.get("lnSpcReduction", "0")) / 100000.0
        return fs, lr
    except Exception:
        return 1.0, 0.0


_STAR5 = "polygon(50% 0%, 61% 35%, 98% 35%, 68% 57%, 79% 91%, 50% 70%, 21% 91%, 32% 57%, 2% 35%, 39% 35%)"


def _geom_css(shape):
    """prstGeom -> CSS so an ellipse isn't a rectangle and a star isn't a square (cycle-3
    pin: 'what were stars have become squares')."""
    try:
        pg = shape._element.spPr.find("{%s}prstGeom" % A)
        prst = pg.get("prst") if pg is not None else None
    except Exception:
        prst = None
    if prst == "ellipse":
        return "border-radius:50%;"
    if prst in ("roundRect", "round2SameRect"):
        return "border-radius:12px;"
    if prst == "star5":
        return "clip-path:%s;" % _STAR5
    if prst == "triangle":
        return "clip-path:polygon(50% 0%, 100% 100%, 0% 100%);"
    if prst == "diamond":
        return "clip-path:polygon(50% 0%, 100% 50%, 50% 100%, 0% 50%);"
    return ""


def _border_css(shape, theme_map):
    """Paint the shape OUTLINE when ln carries a solid fill (cycle-3 pin-8: missing boxes
    around '70% of people…'). noFill ln (e.g. Oracle logo) stays borderless."""
    col, wpt = _line_hex(shape, theme_map)
    if not col:
        return ""
    return "border:%.1fpx solid %s;" % (max(1.0, wpt * 96.0 / 72.0), col)


def _border_css_t(shape, theme_map, recolor):
    css = _border_css(shape, theme_map)
    if css and recolor == "dark":
        m = re.search(r"solid (#[0-9a-fA-F]{6})", css)
        if m and _contrast_ratio(m.group(1), DARK_BG) < 2.2:
            css = css.replace(m.group(1), "#7FB8D9")     # soft sky outline on navy (pin-15)
    return css


def _is_txbox(shape):
    try:
        P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        c = shape._element.find(".//{%s}cNvSpPr" % P)
        return c is not None and c.get("txBox") == "1"
    except Exception:
        return False


_LINE_PRSTS = {"line", "straightConnector1", "bentConnector2", "bentConnector3"}
_SIMPLE_PRSTS = {"rect", "roundRect", "round2SameRect", "ellipse", "star5", "triangle",
                 "diamond", "line", "straightConnector1", "rightArrow", "chevron", "homePlate"}


def _shape_prst(sh):
    try:
        pg = sh._element.spPr.find("{%s}prstGeom" % A)
        return pg.get("prst") if pg is not None else None
    except Exception:
        return None


def is_simple_group(grp):
    """A group whose every descendant is a basic shape/picture/text renders NATIVELY —
    no raster, real clickable text (cycle-3 pin-2). Charts/freeforms/tables stay raster."""
    try:
        for ch in grp.shapes:
            if ch.shape_type == 6:
                if not is_simple_group(ch):
                    return False
                continue
            if ch.shape_type == 13:
                try:
                    ch.image.blob
                except Exception:
                    return False
                continue
            el = ch._element
            if el.tag.endswith("graphicFrame"):
                return False
            if el.find(".//{%s}custGeom" % A) is not None:
                return False
            prst = _shape_prst(ch)
            if prst is not None and prst not in _SIMPLE_PRSTS:
                return False
        return True
    except Exception:
        return False


# ============================================================ geometry transform
def _knockout_white(data_uri, thresh=238):
    """Make the raster's background transparent by FLOOD-FILLING near-white from the
    EDGES ONLY (cycle-7): interior white pixels — the cores of letters, ring labels —
    survive, so text no longer degrades to hollow outlines on dark slides."""
    try:
        import io
        from collections import deque
        from PIL import Image
        raw = base64.b64decode(data_uri.split(",", 1)[1])
        im = Image.open(io.BytesIO(raw)).convert("RGBA")
        px = im.load()
        w, h = im.size
        seen = bytearray(w * h)
        q = deque()
        def is_white(x, y):
            r, g, b, a = px[x, y]
            return a > 0 and r >= thresh and g >= thresh and b >= thresh
        for x in range(w):
            for y in (0, h - 1):
                if is_white(x, y) and not seen[y * w + x]:
                    seen[y * w + x] = 1; q.append((x, y))
        for y in range(h):
            for x in (0, w - 1):
                if is_white(x, y) and not seen[y * w + x]:
                    seen[y * w + x] = 1; q.append((x, y))
        while q:
            x, y = q.popleft()
            r, g, b, a = px[x, y]
            px[x, y] = (r, g, b, 0)
            for nx, ny in ((x-1, y), (x+1, y), (x, y-1), (x, y+1)):
                if 0 <= nx < w and 0 <= ny < h and not seen[ny * w + nx] and is_white(nx, ny):
                    seen[ny * w + nx] = 1; q.append((nx, ny))
        buf = io.BytesIO(); im.save(buf, format="PNG")
        return "data:image/png;base64," + base64.b64encode(buf.getvalue()).decode("ascii")
    except Exception:
        return data_uri


def _canvas_tf(sx, sy):
    """Top-level shape space (EMU) -> canvas px. Returns fn(l,t,w,h)->(x,y,w,h)."""
    return lambda l, t, w, h: (l * sx, t * sy, w * sx, h * sy)


def _group_tf(parent_tf, group_el):
    """Compose a child-space transform for a GroupShape from its <a:xfrm> off/ext/chOff/chExt."""
    try:
        # grpSpPr is in the PRESENTATIONML namespace — the a:-namespaced find NEVER matched,
        # so group children silently rendered at RAW child coordinates (cycle-3: exploded
        # number circles). Third instance of this bug family (spPr fill, spPr ln, grpSpPr).
        P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        gp = group_el.find("{%s}grpSpPr" % P)
        xf = gp.find("{%s}xfrm" % A) if gp is not None else None
        off = xf.find("{%s}off" % A); ext = xf.find("{%s}ext" % A)
        cho = xf.find("{%s}chOff" % A); che = xf.find("{%s}chExt" % A)
        ox, oy = int(off.get("x")), int(off.get("y"))
        ex, ey = int(ext.get("cx")), int(ext.get("cy"))
        cx, cy = int(cho.get("x")), int(cho.get("y"))
        ccx, ccy = int(che.get("cx")), int(che.get("cy"))
        sxx = ex / ccx if ccx else 1.0
        syy = ey / ccy if ccy else 1.0
    except Exception:
        return parent_tf
    def tf(l, t, w, h):
        pl = ox + (l - cx) * sxx
        pt = oy + (t - cy) * syy
        return parent_tf(pl, pt, w * sxx, h * syy)
    # cumulative FONT scale: PowerPoint scales glyphs inside transformed groups
    tf.fscale = getattr(parent_tf, "fscale", 1.0) * ((sxx + syy) / 2.0)
    return tf


def _pos_style(x, y, w, h):
    st = "position:absolute;left:%.1fpx;top:%.1fpx;" % (x, y)
    if w:
        st += "width:%.1fpx;" % w
    if h:
        st += "height:%.1fpx;" % h
    return st


# ============================================================ fill / gradient / line
def _fill_alpha(sh):
    """<a:alpha> on the shape's solid fill (0..1). Solid fills were painted opaque —
    review pin-18: a 60%-alpha black scrim rendered as solid black."""
    try:
        # NOTE: spPr is in the PRESENTATIONML namespace (p:), not drawingml — use the
        # python-pptx property (fixed a bug where alpha was never read: pin-18)
        spPr = getattr(sh._element, "spPr", None)
        sf = spPr.find("{%s}solidFill" % A) if spPr is not None else None
        if sf is None:
            return 1.0
        a = sf.find(".//{%s}alpha" % A)
        return int(a.get("val")) / 100000.0 if a is not None else 1.0
    except Exception:
        return 1.0


def _matches_slide_bg(fill):
    try:
        if not fill or fill.startswith("rgba"):
            return False
        r1, g1, b1 = _rgb(fill); r2, g2, b2 = _rgb(_CUR_SLIDE_BG[0])
        return abs(r1 - r2) + abs(g1 - g2) + abs(b1 - b2) <= 18
    except Exception:
        return False


def _fill_hex(sh, theme_map):
    """Shape fill css colour. XML-first so colour MODIFIERS (lumMod/tint/shade) and alpha
    apply — python-pptx fore_color ignores them (cycle-1: tinted panels rendered fully
    saturated). Falls back to the python-pptx path."""
    try:
        spPr = getattr(sh._element, "spPr", None)
        if spPr is not None:
            css, _hex = _solidfill_color(spPr.find("{%s}solidFill" % A), theme_map)
            if css:
                return css
        f = sh.fill
        if f.type == 1:
            fc = f.fore_color
            if fc.type == 1:
                return "#" + str(fc.rgb)
            name = getattr(fc.theme_color, "name", None)
            key = _THEME_KEY.get(name or "")
            if key and theme_map.get(key):
                return "#" + theme_map[key]
    except Exception:
        pass
    return None


def _gradient_css(sh, theme_map, recolor=None):
    try:
        grad = sh._element.find(".//{%s}gradFill" % A)
        if grad is None:
            return None
        stops = []
        for gs in grad.findall(".//{%s}gs" % A):
            pos = int(gs.get("pos", "0")) / 100000.0
            srgb = gs.find("{%s}srgbClr" % A); sch = gs.find("{%s}schemeClr" % A)
            clr = srgb if srgb is not None else sch
            if clr is None:
                continue
            hexc = srgb.get("val") if srgb is not None else _scheme_hex(sch.get("val", ""), theme_map)
            if not hexc:
                continue
            if recolor:
                hexc = _map_color("#" + hexc, recolor).lstrip("#")
            alpha = 1.0
            a = clr.find("{%s}alpha" % A)
            if a is not None:
                alpha = int(a.get("val")) / 100000.0
            r, g, b = _rgb("#" + hexc)
            stops.append((pos, "rgba(%d,%d,%d,%.3f)" % (r, g, b, alpha)))
        if not stops:
            return None
        stops.sort()
        lin = grad.find("{%s}lin" % A)
        ang = int(lin.get("ang")) / 60000.0 if (lin is not None and lin.get("ang")) else 0.0
        css_ang = (90 + ang) % 360
        return "linear-gradient(%.1fdeg, %s)" % (css_ang, ", ".join("%s %.1f%%" % (c, p * 100) for p, c in stops))
    except Exception:
        return None


_LN_STYLE_W = None  # cache of theme lnStyleLst widths (pt)


def _theme_ln_widths(sh):
    global _LN_STYLE_W
    if _LN_STYLE_W is not None:
        return _LN_STYLE_W
    try:
        theme = sh.part.slide_layout.slide_master.element.getroottree()  # unused fallback
    except Exception:
        pass
    _LN_STYLE_W = [1.0, 2.0, 3.0]   # OOXML default-ish: subtle/moderate/intense (pt)
    return _LN_STYLE_W


def _style_ln(sh, theme_map):
    """(color, width_pt) from <p:style><a:lnRef> — the theme-styled outline that shapes
    use when they carry no explicit ln fill (cycle-4 pins 7/8/9)."""
    try:
        P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        st = sh._element.find("{%s}style" % P)
        if st is None:
            return None, 0.0
        lnref = st.find("{%s}lnRef" % A)
        if lnref is None or int(lnref.get("idx", "0")) == 0:
            return None, 0.0
        sch = lnref.find("{%s}schemeClr" % A)
        if sch is None:
            return None, 0.0
        hexc = _scheme_hex(sch.get("val", ""), theme_map)
        if not hexc:
            return None, 0.0
        col = _apply_mods("#" + hexc, sch)
        w = _theme_ln_widths(sh)[min(2, int(lnref.get("idx")) - 1)]
        return col, w
    except Exception:
        return None, 0.0


def _line_hex(sh, theme_map):
    """Connector / shape outline colour (for LINE shapes and borders)."""
    try:
        # spPr lives in the p: namespace — the old a:-namespaced find NEVER matched, so
        # connector/divider lines silently vanished (review pin-22)
        spPr = getattr(sh._element, "spPr", None)
        ln = spPr.find("{%s}ln" % A) if spPr is not None else None
        if ln is not None and ln.find("{%s}noFill" % A) is not None:
            return None, 0.0                              # explicit noFill NEVER draws (Oracle pin)
        if ln is None or ln.find("{%s}solidFill" % A) is None:
            return _style_ln(sh, theme_map)               # theme styleRef outline (cycle-4)
        if ln is None:
            return None, 1.0
        w = ln.get("w")
        width = (int(w) / EMU_PER_PT) if w else 1.0
        sf = ln.find("{%s}solidFill" % A)
        if sf is None:
            return None, width
        srgb = sf.find("{%s}srgbClr" % A); sch = sf.find("{%s}schemeClr" % A)
        hexc = srgb.get("val") if srgb is not None else _scheme_hex(sch.get("val", ""), theme_map)
        return ("#" + hexc if hexc else None), width
    except Exception:
        return None, 1.0


# ============================================================ text (RC4)
def _typeface(tf, font_scheme):
    if not tf:
        return None
    if tf in ("+mj-lt", "+mj-ea", "+mj-cs"):
        return font_scheme.get("major", "Montserrat")
    if tf in ("+mn-lt", "+mn-ea", "+mn-cs"):
        return font_scheme.get("minor", "Montserrat")
    return tf


def _props_from_rpr(el, theme_map, font_scheme):
    """Extract {bold,italic,size(pt),font,color} from an rPr/defRPr/endParaRPr element."""
    d = {}
    if el is None:
        return d
    b, i, sz = el.get("b"), el.get("i"), el.get("sz")
    if b is not None:
        d["bold"] = b in ("1", "true")
    if i is not None:
        d["italic"] = i in ("1", "true")
    if sz is not None:
        try:
            d["size"] = int(sz) / 100.0
        except Exception:
            pass
    lat = el.find("{%s}latin" % A)
    if lat is not None:
        f = _typeface(lat.get("typeface"), font_scheme)
        if f:
            d["font"] = f
    sf = el.find("{%s}solidFill" % A)
    if sf is not None:
        css, hexc = _solidfill_color(sf, theme_map)   # modifiers + alpha applied
        if hexc:
            d["color"] = hexc
    return d


def _placeholder_defaults(shape, theme_map, font_scheme):
    """Per-level run defaults inherited from the LAYOUT placeholder, MASTER placeholder and
    the master's txStyles — where a Title's 44pt/colour actually live (loop cycle-1: s01's
    title rendered ~11px because only the shape's own lstStyle was consulted).
    Returns {level(0-based): props}. Merge order (weakest first): master txStyles ->
    master placeholder lstStyle -> layout placeholder lstStyle."""
    out = {}
    try:
        if not shape.is_placeholder:
            return out
        ph = shape.placeholder_format
        _PH_XML = {"TITLE": "title", "CENTER_TITLE": "ctrTitle", "SUBTITLE": "subTitle",
                   "BODY": "body", "SLIDE_NUMBER": "sldNum", "FOOTER": "ftr", "DATE": "dt",
                   "OBJECT": "body", "PICTURE": "pic", "CHART": "chart", "TABLE": "tbl"}
        ph_enum = str(ph.type).split(".")[-1].split(" ")[0] if ph.type is not None else ""
        ph_type = _PH_XML.get(ph_enum, ph_enum.lower() or "body")
        ph_idx = ph.idx
        slide = shape.part.package  # placeholder; real slide fetched via _parent chain below
    except Exception:
        return out

    def lst_levels(lst_el):
        levels = {}
        if lst_el is None:
            return levels
        for n in range(1, 10):
            lvl = lst_el.find("{%s}lvl%dpPr" % (A, n))
            if lvl is not None:
                dr = lvl.find("{%s}defRPr" % A)
                pr = _props_from_rpr(dr, theme_map, font_scheme) if dr is not None else {}
                if lvl.get("algn"):
                    pr["align"] = {"ctr": "center", "r": "right", "just": "justify"}.get(lvl.get("algn"), "left")
                if lvl.get("marL"):
                    pr["marL"] = int(lvl.get("marL"))
                if lvl.get("indent"):
                    pr["indent"] = int(lvl.get("indent"))
                bu = lvl.find("{%s}buChar" % A)
                if bu is not None:
                    pr["bullet"] = _map_bullet_char(bu.get("char", ""), lvl)
                if lvl.find("{%s}buNone" % A) is not None:
                    pr["bullet"] = ""
                if pr:
                    levels[n - 1] = pr
        return levels

    def merge(dst, src):
        for k, v in src.items():
            base = dict(dst.get(k, {})); base.update(v); dst[k] = base

    def find_ph(shapes_el):
        """Matching placeholder <p:sp> on a layout/master by idx, else by type."""
        P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        for sp in shapes_el.iter("{%s}sp" % P):
            phe = sp.find(".//{%s}ph" % P)
            if phe is None:
                continue
            if ph_idx is not None and phe.get("idx") and int(phe.get("idx")) == ph_idx:
                return sp
        for sp in shapes_el.iter("{%s}sp" % P):
            phe = sp.find(".//{%s}ph" % P)
            if phe is not None and phe.get("type", "body") == ph_type:
                return sp
        return None

    try:
        slide_obj = shape._parent
        while slide_obj is not None and not hasattr(slide_obj, "slide_layout"):
            slide_obj = getattr(slide_obj, "_parent", None)
        layout = slide_obj.slide_layout if slide_obj is not None else None
        master = layout.slide_master if layout is not None else None
    except Exception:
        return out

    # 1) master txStyles (titleStyle for titles, bodyStyle otherwise)
    try:
        P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        tx = master.element.find("{%s}txStyles" % P)
        if tx is not None:
            if "title" in ph_type.lower():
                style_tag = "titleStyle"
            elif ph_type in ("sldNum", "ftr", "dt"):
                style_tag = "otherStyle"
            else:
                style_tag = "bodyStyle"
            merge(out, lst_levels(tx.find("{%s}%s" % (P, style_tag))))
    except Exception:
        pass
    # 2) master placeholder, 3) layout placeholder (stronger)
    for host in (master, layout):
        try:
            sp = find_ph(host.element)
            if sp is not None:
                txb = sp.find(".//{%s}txBody" % "http://schemas.openxmlformats.org/presentationml/2006/main")
                if txb is None:
                    txb = sp.find(".//{%s}txBody" % A)
                lst = txb.find("{%s}lstStyle" % A) if txb is not None else None
                merge(out, lst_levels(lst))
        except Exception:
            pass
    return out


def _shape_defaults(shape, theme_map, font_scheme):
    """Shape-level lstStyle lvl1 defRPr — where a Title's bold/size/colour often live (RC4)."""
    try:
        txb = shape.text_frame._txBody
        lst = txb.find("{%s}lstStyle" % A)
        if lst is None:
            return {}
        lvl = lst.find("{%s}lvl1pPr" % A)
        dr = lvl.find("{%s}defRPr" % A) if lvl is not None else None
        pr = _props_from_rpr(dr, theme_map, font_scheme)
        if lvl is not None and lvl.get("algn"):
            pr["align"] = {"ctr": "center", "r": "right", "just": "justify"}.get(lvl.get("algn"), "left")
        return pr
    except Exception:
        return {}


def _iter_para(para, shape_def, theme_map, font_scheme, ph_defs=None):
    """Yield ('run', text, props) and ('br',) tokens for a paragraph, resolving each run's
    props through: layout/master placeholder defaults (per paragraph LEVEL) -> shape
    lstStyle defRPr -> paragraph pPr/defRPr -> run rPr (run wins).
    Iterating the XML (not python-pptx .runs) preserves <a:br> soft line-breaks (RC4)."""
    p = para._p
    ppr = p.find("{%s}pPr" % A)
    lvl = 0
    if ppr is not None and ppr.get("lvl"):
        try:
            lvl = int(ppr.get("lvl"))
        except Exception:
            pass
    para_def = _props_from_rpr(ppr.find("{%s}defRPr" % A) if ppr is not None else None, theme_map, font_scheme)
    base = dict((ph_defs or {}).get(lvl, {}))
    base.update(shape_def); base.update(para_def)
    for child in p:
        tag = child.tag.split("}")[1]
        if tag == "r":
            t = "".join(e.text or "" for e in child.findall("{%s}t" % A))
            props = dict(base)
            props.update(_props_from_rpr(child.find("{%s}rPr" % A), theme_map, font_scheme))
            yield ("run", t, props)
        elif tag == "br":
            yield ("br", "", {})
        elif tag == "fld":
            t = "".join(e.text or "" for e in child.findall("{%s}t" % A))
            if t:
                props = dict(base)
                props.update(_props_from_rpr(child.find("{%s}rPr" % A), theme_map, font_scheme))
                yield ("run", t, props)


_WINGDINGS = {"§": "▪", "Ø": "➢", "ü": "✓",
              "v": "❖", "l": "●", "n": "■", "u": "◆"}


def _map_bullet_char(ch, lvl_el):
    """buChar with a symbol font maps to unicode so bullets don't render as tofu boxes."""
    try:
        bf = lvl_el.find("{%s}buFont" % A)
        if bf is not None and (bf.get("typeface") or "").startswith(("Wingdings", "Webdings")):
            return _WINGDINGS.get(ch, "•")
    except Exception:
        pass
    return ch


def _para_bullet(ppr, base):
    """(bullet_char, marL_emu, indent_emu) for a paragraph: explicit pPr wins, else the
    inherited level defaults. 235 buChar bullets in the collection rendered as NOTHING
    before this (cycle-5 finding)."""
    bullet = base.get("bullet", "")
    marL = base.get("marL", 0)
    indent = base.get("indent", 0)
    if ppr is not None:
        if ppr.get("marL"):
            marL = int(ppr.get("marL"))
        if ppr.get("indent"):
            indent = int(ppr.get("indent"))
        bu = ppr.find("{%s}buChar" % A)
        if bu is not None:
            bullet = _map_bullet_char(bu.get("char", ""), ppr)
        if ppr.find("{%s}buNone" % A) is not None:
            bullet = ""
    if bullet and not marL:
        marL = 228600                      # default hang when a bullet exists
        indent = indent or -228600
    return bullet, marL, indent


def _ph_bodypr(shape):
    """Merged bodyPr attributes (anchor, lIns/tIns/rIns/bIns) this placeholder INHERITS from
    the layout/master placeholder. bodyPr attributes cascade per-attribute (master weakest,
    then layout), so a layout title with a bare <a:bodyPr/> still inherits the master's
    anchor='ctr' (cycle-10 pin-3). Matched by idx, else by type. Returns a dict (may be
    empty); attribute absent everywhere => key absent."""
    out = {}
    try:
        ph = shape.placeholder_format
        ph_idx = ph.idx
        _PH_XML = {"TITLE": "title", "CENTER_TITLE": "ctrTitle", "SUBTITLE": "subTitle",
                   "BODY": "body", "SLIDE_NUMBER": "sldNum", "FOOTER": "ftr", "DATE": "dt",
                   "OBJECT": "body", "PICTURE": "pic", "CHART": "chart", "TABLE": "tbl"}
        ph_enum = str(ph.type).split(".")[-1].split(" ")[0] if ph.type is not None else ""
        ph_type = _PH_XML.get(ph_enum, ph_enum.lower() or "body")
        slide_obj = shape._parent
        while slide_obj is not None and not hasattr(slide_obj, "slide_layout"):
            slide_obj = getattr(slide_obj, "_parent", None)
        layout = slide_obj.slide_layout if slide_obj is not None else None
        master = layout.slide_master if layout is not None else None
        P = "http://schemas.openxmlformats.org/presentationml/2006/main"

        def find_bp(shapes_el):
            hit = None
            for sp in shapes_el.iter("{%s}sp" % P):     # idx match preferred
                phe = sp.find(".//{%s}ph" % P)
                if phe is not None and ph_idx is not None and phe.get("idx") and int(phe.get("idx")) == ph_idx:
                    hit = sp; break
            if hit is None:
                for sp in shapes_el.iter("{%s}sp" % P):  # else type match
                    phe = sp.find(".//{%s}ph" % P)
                    if phe is not None and phe.get("type", "body") == ph_type:
                        hit = sp; break
            return hit.find(".//{%s}bodyPr" % A) if hit is not None else None

        for host in (master, layout):                   # master weakest, layout overrides
            if host is None:
                continue
            bp = find_bp(host.element)
            if bp is not None:
                for k in ("anchor", "lIns", "tIns", "rIns", "bIns"):
                    if bp.get(k) is not None:
                        out[k] = bp.get(k)
    except Exception:
        pass
    return out


def _render_text_shape(shape, x, y, w, h, sx, theme_map, font_scheme, recolor, bg):
    box = _pos_style(x, y, w, h)
    # bodyPr: vertical anchor (PowerPoint DEFAULT is TOP — centring everything was a
    # cycle-1 defect) + text insets (lIns/tIns/rIns/bIns, defaults 91440/45720 EMU)
    # drawn autoshapes default to MIDDLE anchor in PowerPoint; text boxes/placeholders to top
    anchor = "flex-start" if (_is_txbox(shape) or getattr(shape, "is_placeholder", False)) else "center"
    try:
        # tiny single-paragraph boxes (decorative digits over discs) read centred in PPT
        if len(shape.text_frame.paragraphs) == 1 and h <= 90 and len(shape.text_frame.text.strip()) <= 4:
            anchor = "center"
    except Exception:
        pass
    l_in, t_in, r_in, b_in = 91440, 45720, 91440, 45720
    try:
        bp = shape.text_frame._txBody.find("{%s}bodyPr" % A)
        # A placeholder's bodyPr attrs (anchor + insets) INHERIT from the layout/master
        # placeholder when absent on the slide shape — a bare <a:bodyPr/> on a title
        # inherits anchor="ctr" from the master (cycle-10 pin-3: the title rendered
        # top-anchored, overflowing DOWN onto the chart, when PowerPoint centres it).
        ibp = _ph_bodypr(shape) if getattr(shape, "is_placeholder", False) else {}

        def _attr(name):
            if bp is not None and bp.get(name) is not None:
                return bp.get(name)
            return ibp.get(name)
        anch = _attr("anchor")
        if anch:
            anchor = {"ctr": "center", "b": "flex-end", "t": "flex-start"}.get(anch, anchor)
        l_in = int(_attr("lIns") or l_in); t_in = int(_attr("tIns") or t_in)
        r_in = int(_attr("rIns") or r_in); b_in = int(_attr("bIns") or b_in)
    except Exception:
        pass
    box += ("display:flex;flex-direction:column;justify-content:%s;overflow:visible;box-sizing:border-box;"
            "padding:%.1fpx %.1fpx %.1fpx %.1fpx;" % (anchor, t_in * sx, r_in * sx, b_in * sx, l_in * sx))
    if bg:
        box += "background:%s;" % bg
    box += _geom_css(shape) + _border_css_t(shape, theme_map, recolor)
    fs_scale, ln_red = _autofit_scale(shape)
    fit_w = None
    fit_shrink = 1.0
    try:
        bp0 = shape.text_frame._txBody.find("{%s}bodyPr" % A)
        is_title_ph = getattr(shape, "is_placeholder", False) and "title" in str(
            getattr(shape.placeholder_format, "type", "")).lower()
        # cycle-10 pin-3: an EXPLICIT fontScale is PowerPoint's own computed fit —
        # authoritative. PPT renders sz*fontScale and lets the text overflow the box
        # (s05 title). Estimating our own shrink on top double-shrank (90% -> 61%),
        # so the height estimator only runs when PPT stored no scale.
        if fs_scale >= 1.0 and ((bp0 is not None and bp0.find("{%s}normAutofit" % A) is not None) or is_title_ph):
            fit_w = max(20.0, w - (l_in + r_in) * sx)
            usable_h = max(20.0, h - (t_in + b_in) * sx)
            # estimate total height at declared sizes; shrink uniformly if it overflows
            est_h = 0.0
            _pd = _placeholder_defaults(shape, theme_map, font_scheme)
            _sd = _shape_defaults(shape, theme_map, font_scheme)
            for _para in shape.text_frame.paragraphs:
                _txt = (_para.text or "").strip()
                _sz = None
                for _tok in _iter_para(_para, _sd, theme_map, font_scheme, _pd):
                    if _tok[0] == "run" and _tok[2].get("size"):
                        _sz = _tok[2]["size"]; break
                _szpx = (_sz or 18.0) * fs_scale * EMU_PER_PT * sx
                if not _txt:
                    est_h += _szpx * 1.16; continue
                _avg = 0.60 if _sd.get("bold") or (_pd.get(0, {}) or {}).get("bold") else 0.55
                _lines = max(1, int((len(_txt) * _avg * _szpx) / fit_w + 0.995))
                est_h += _lines * _szpx * 1.16
            if est_h > usable_h * 1.04:
                fit_shrink = max(0.68, usable_h / est_h)
    except Exception:
        fit_w, fit_shrink = None, 1.0
    shape_def = _shape_defaults(shape, theme_map, font_scheme)
    ph_defs = _placeholder_defaults(shape, theme_map, font_scheme)
    inner = []
    _ALGN_XML = {"ctr": "center", "r": "right", "just": "justify", "l": "left"}
    for para in shape.text_frame.paragraphs:
        align = None
        if para.alignment is not None:
            align = _ALIGN.get(str(para.alignment).split(".")[-1].split(" ")[0], None)
        if align is None:                       # inherit: shape lstStyle -> placeholder lvl
            lvl0 = 0
            _ppr = para._p.find("{%s}pPr" % A)
            if _ppr is not None and _ppr.get("lvl"):
                try: lvl0 = int(_ppr.get("lvl"))
                except Exception: pass
            align = shape_def.get("align") or (ph_defs.get(lvl0, {}) or {}).get("align") or "left"
        line_spans, has_text = [], False
        for tok in _iter_para(para, shape_def, theme_map, font_scheme, ph_defs):
            if tok[0] == "br":
                line_spans.append("<br>")
                continue
            _, t, props = tok
            if not t:
                continue
            has_text = True
            s = ["letter-spacing:-0.012em", "line-height:1.16"]
            if props.get("size"):
                fpx = props["size"] * fs_scale * fit_shrink * EMU_PER_PT * sx
                if fit_w:
                    txt_len = max(1, len((para.text or "").strip()))
                    avg_em = 0.60 if props.get("bold") else 0.55
                    need = fit_w / (avg_em * txt_len)
                    if need < fpx:
                        fpx = max(need, fpx * 0.72)       # shrink like PPT, floor at 72%
                s.append("font-size:%.1fpx" % fpx)
            if ln_red:
                s.append("line-height:%.2f" % max(0.9, 1.2 * (1 - ln_red)))
            s.append("font-weight:%d" % (700 if props.get("bold") else 400))
            if props.get("italic"):
                s.append("font-style:italic")
            fam = props.get("font") or font_scheme.get("minor", "Montserrat")
            s.append("font-family:'%s',Montserrat,sans-serif" % fam)
            col = props.get("color")
            if recolor:
                col = _map_text_color(col, recolor)
            if recolor == "dark":
                _is_title = getattr(shape, "is_placeholder", False) and "title" in str(
                    getattr(shape.placeholder_format, "type", "")).lower()
                if _is_title or (props.get("size") or 0) >= 26:
                    col = "#F2F7FC"                     # cycle-8: titles are ALWAYS white on dark
            if col:
                s.append("color:%s" % col)
            line_spans.append('<span style="%s">%s</span>' % (";".join(s), _html.escape(t)))
        if has_text and line_spans and line_spans[-1] == "<br>":
            # cycle-10 pin-14: PPT renders a TRAILING soft break (<a:br> before
            # endParaRPr, e.g. 'Power User\x0b') as a visible empty line; HTML
            # collapses a trailing <br>, so the column's bullets started a line
            # high. Materialize the line box with an nbsp styled like the last run.
            _last = next((sp for sp in reversed(line_spans) if sp.startswith("<span")), "")
            _mst = re.match(r'<span style="([^"]*)"', _last)
            line_spans.append('<span style="%s">&#160;</span>' % (_mst.group(1) if _mst else ""))
        _ppr2 = para._p.find("{%s}pPr" % A)
        _spc_css = ""
        if _ppr2 is not None:
            for tag, css in (("spcBef", "margin-top"), ("spcAft", "margin-bottom")):
                el2 = _ppr2.find("{%s}%s/{%s}spcPts" % (A, tag, A))
                if el2 is not None and el2.get("val"):
                    _spc_css += "%s:%.1fpx;" % (css, int(el2.get("val")) / 100.0 * EMU_PER_PT * sx)
            _ln2 = _ppr2.find("{%s}lnSpc/{%s}spcPct" % (A, A))
            if _ln2 is not None and _ln2.get("val"):
                _spc_css += "line-height:%.2f;" % (int(_ln2.get("val")) / 100000.0 * 1.2)
        _lvl2 = 0
        if _ppr2 is not None and _ppr2.get("lvl"):
            try: _lvl2 = int(_ppr2.get("lvl"))
            except Exception: pass
        _base2 = dict((ph_defs or {}).get(_lvl2, {})); _base2.update(shape_def)
        bullet, marL, indent = _para_bullet(_ppr2, _base2)
        if para._p.find("{%s}fld" % A) is not None:
            bullet = ""                       # slide-number/date fields never carry bullets
        pstyle = "text-align:%s;" % align + _spc_css
        if marL:
            pstyle += "padding-left:%.1fpx;" % (marL * sx)
        if indent:
            pstyle += "text-indent:%.1fpx;" % (indent * sx)
        if has_text or line_spans:
            _bw = abs(indent) * sx if indent else 24.0
            prefix = ('<span style="display:inline-block;width:%.1fpx;text-indent:0">%s</span>'
                      % (_bw, _html.escape(bullet))) if (bullet and has_text) else ""
            inner.append('<div style="%s">%s</div>' % (pstyle, prefix + ("".join(line_spans) or "&nbsp;")))
        else:
            _eps = para._p.find("{%s}endParaRPr" % A)
            _esz = None
            if _eps is not None and _eps.get("sz"):
                try: _esz = int(_eps.get("sz")) / 100.0
                except Exception: pass
            _esz = _esz or _base2.get("size") or 12.0
            inner.append('<div style="font-size:%.1fpx;line-height:1.16">&nbsp;</div>' % (_esz * fs_scale * EMU_PER_PT * sx))
    return '<div style="%s">%s</div>' % (box, "".join(inner))


# ============================================================ shape dispatch
def _render_shape(shape, tf, sx, theme_map, font_scheme, recolor=None, swap_image=False, target=None):
    """One shape -> HTML at its transformed canvas box. Handles picture (crop/stretch),
    connector line, gradient/solid fill, and text; recolour applied when in recolor mode."""
    try:
        l, t, w, h = shape.left, shape.top, shape.width, shape.height
    except Exception:
        return ""
    if l is None or t is None:
        return ""
    x, y, ww, hh = tf(l, t, w or 0, h or 0)

    st = shape.shape_type
    _prst = _shape_prst(shape)
    # picture — RC2 crop/stretch semantics (incl. negative srcRect) via _pic_html
    if st == 13:
        try:
            img = shape.image
            fx = _blip_fx(shape, img.blob, theme_map, recolor)
            blob, ctype = (fx, "image/png") if fx else (img.blob, img.content_type)
            uri = _process_image(blob, ctype, target) if swap_image else None
            if uri is None:
                uri = _embed_image(blob, ctype)
            return _pic_html(uri, x, y, ww, hh, _src_rect(shape))
        except Exception as e:
            # review pin-12: a picture must NEVER be silently dropped. Fall back to a raw
            # blob embed (browsers render svg/most types from data URIs); warn loudly so
            # a truly unrenderable format (emf/wmf -> needs COM raster) is visible.
            try:
                part = shape._element.blipFill
                rId = part.blip.rEmbed
                image_part = shape.part.related_part(rId)
                ct = image_part.content_type or "image/png"
                uri = "data:%s;base64,%s" % (ct, base64.b64encode(image_part.blob).decode("ascii"))
                print("  [warn] picture %r fell back to raw embed (%s): %s" % (getattr(shape, "name", "?"), ct, e))
                return '<div style="%sbackground-image:url(\'%s\');background-size:100%% 100%%;background-repeat:no-repeat;"></div>' % (
                    _pos_style(x, y, ww, hh), uri)
            except Exception as e2:
                print("  [WARN] picture %r DROPPED (unrenderable: %s / %s) — needs COM raster" % (getattr(shape, "name", "?"), e, e2))
                return ""
    # connector / line (RC3) — incl. line-GEOMETRY autoshapes (chart axes, separators:
    # cycle-3 pins 9/13) — thin rect along the box using the line colour
    if st == 9 or _prst in _LINE_PRSTS:
        col, lw = _line_hex(shape, theme_map)
        if col:
            if recolor:
                col = _map_color(col, recolor)
            thick = max(1.0, lw * sx * EMU_PER_PT / EMU_PER_PT)  # lw already in pt
            thick = max(1.0, lw)
            if hh >= ww:   # vertical
                return '<div style="%sbackground:%s;"></div>' % (_pos_style(x, y, max(thick, 1), hh), col)
            return '<div style="%sbackground:%s;"></div>' % (_pos_style(x, y, ww, max(thick, 1)), col)
        return ""
    # gradient-only decorative shape (the 'fade' scrim), no text
    grad = _gradient_css(shape, theme_map, recolor)
    has_txt = shape.has_text_frame and shape.text_frame.text.strip()
    if grad and not has_txt:
        return '<div style="%sbackground:%s;"></div>' % (_pos_style(x, y, ww, hh), grad)
    # text
    if has_txt:
        bg = grad or _fill_hex(shape, theme_map)
        if bg and not grad and _matches_slide_bg(bg):
            bg = None                                   # invisible-on-source fill stays invisible
        if bg and recolor and not grad:
            bg = _map_color(bg, recolor)
        return _render_text_shape(shape, x, y, ww, hh, sx * getattr(tf, "fscale", 1.0), theme_map, font_scheme, recolor, bg)
    # solid-filled / outlined autoshape (no text) — with geometry + border
    fill = _fill_hex(shape, theme_map)
    extra = _geom_css(shape) + _border_css_t(shape, theme_map, recolor)
    if fill or extra.startswith("border") or "border:" in extra:
        if fill and recolor:
            fill = _map_color(fill, recolor)
        return '<div style="%s%s%s"></div>' % (_pos_style(x, y, ww, hh),
                                               ("background:%s;" % fill) if fill else "", extra)
    return ""


def _walk_shapes(shapes, tf, sx, theme_map, font_scheme, recolor, flip, prs, raster_shapes, top=True,
                 suppress_names=None, flip_target=None, logo_swaps=None, dark_backers=None,
                 dark_display=False):
    """Render a shape collection. Top-level complex shapes with a COM raster are emitted as
    an image; groups without a raster are recursed (best-effort); everything else dispatched.
    suppress_names: curation-driven per-slide shape suppression (review pin-2 — e.g. a
    decorative line the customer ruled out; listed in curation.json shape_suppressions).
    dark_backers: shape names that, when the DISPLAYED theme is dark, get a soft light chip
    drawn immediately BEHIND them (cycle-10 pin-6: dark competitor logos vanished on the
    dark card). Emitted just before the shape so it paints above the card, below the logo."""
    parts = []
    for i, sh in enumerate(shapes):
        if suppress_names and getattr(sh, "name", None) in suppress_names:
            continue
        if dark_backers and dark_display and getattr(sh, "name", None) in dark_backers:
            try:
                bl, bt, bw, bh = tf(sh.left, sh.top, sh.width or 0, sh.height or 0)
                pad = 12.0
                parts.append('<div style="%sbackground:#EEF3F8;border-radius:12px;'
                             'box-shadow:0 2px 10px rgba(0,0,0,0.28);"></div>'
                             % _pos_style(bl - pad, bt - pad, bw + 2 * pad, bh + 2 * pad))
            except Exception:
                pass
        if logo_swaps and getattr(sh, "name", None) in logo_swaps and sh.shape_type == 13:
            try:                                        # real brand SVG instead of a picture
                l2, t2, w2, h2 = tf(sh.left, sh.top, sh.width or 0, sh.height or 0)
                svg = logo_swaps[sh.name]
                parts.append('<div style="%s">%s</div>' % (_pos_style(l2, t2, w2, h2), svg))
                continue
            except Exception:
                pass
        if top and raster_shapes and i in raster_shapes:
            try:
                l, t, w, h = sh.left, sh.top, sh.width, sh.height
                x, y, ww, hh = tf(l, t, w or 0, h or 0)
                uri = raster_shapes[i]
                uri = _knockout_white(uri)   # white bg transparent: underlying fills show through
                parts.append('<img src="%s" style="%sobject-fit:contain;" alt="">'
                             % (uri, _pos_style(x, y, ww, hh)))
                continue
            except Exception:
                pass
        if sh.shape_type == 6:   # GROUP — recurse with composed transform (RC3 tier 1)
            try:
                child_tf = _group_tf(tf, sh._element)
                parts.append(_walk_shapes(sh.shapes, child_tf, sx, theme_map, font_scheme,
                                          recolor, flip, prs, None, top=False, flip_target=flip_target))
                continue
            except Exception:
                pass
        fb = top and _is_fullbleed(sh, prs)
        parts.append(_render_shape(sh, tf, sx, theme_map, font_scheme,
                                   recolor=recolor, swap_image=(flip and fb),
                                   target=(recolor or flip_target)))
    return "\n".join(p for p in parts if p)


def _is_fullbleed(sh, prs):
    try:
        return bool(sh.width and sh.height and sh.left is not None
                    and sh.width >= prs.slide_width * 0.9 and sh.height >= prs.slide_height * 0.9)
    except Exception:
        return False


_CUR_SLIDE_BG = ["#FFFFFF"]


def render_slide(prs, index0, theme_map, font_scheme=None, theme=None, with_background=True, raster_shapes=None, suppress_names=None, logo_swaps=None, dark_backers=None):
    """Render slide index0 -> `.stage` inner HTML fragment.

    theme=None renders NATIVE (fully faithful). theme='light'|'dark': if it matches the
    native theme the render is faithful; otherwise the ONE sanctioned transform runs —
    colour conversion to the brand guide (RC1). For image-template slides a full-bleed
    background image is brightened/darkened instead. raster_shapes maps top-level shape
    indices to COM-exported PNG data URIs for shapes that can't be rebuilt (RC3)."""
    font_scheme = font_scheme or {"major": "Montserrat", "minor": "Montserrat"}
    sx, sy = _sx(prs), _sy(prs)
    slide = prs.slides[index0]
    native = native_theme_of(prs, slide, theme_map)
    target = theme or native
    flip = target != native

    has_image = any(sh.shape_type == 13 and _is_fullbleed(sh, prs) for sh in slide.shapes)
    if not flip:
        mode, recolor = "faithful", None
    elif has_image:
        mode, recolor = "image-swap", None       # image templates keep style, swap the photo
        img_target = target
    else:
        mode, recolor = "recolor", target         # solid-bg content -> recolour to brand guide

    parts = []
    try:
        _CUR_SLIDE_BG[0] = slide_background(slide, theme_map)
    except Exception:
        _CUR_SLIDE_BG[0] = "#FFFFFF"
    if with_background:
        if mode == "recolor":
            bg = LIGHT_BG if target == "light" else DARK_BG
        else:
            bg = slide_background(slide, theme_map)
            if flip:                              # image-swap: recolour any solid page bg too
                bg = _map_color(bg, target)
        parts.append('<div style="position:absolute;inset:0;background:%s;z-index:0;"></div>' % bg)

    tf = _canvas_tf(sx, sy)
    # LAYOUT/MASTER decoration (cycle-1 pin-1: s01's background image lives on the layout,
    # never the slide). Render NON-placeholder layout+master shapes UNDER slide shapes —
    # placeholders are prompts, not content. Honour showMasterSp="0".
    try:
        layout = slide.slide_layout
        P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        show_master = layout.element.get("showMasterSp", "1") != "0"
        hosts = ([layout.slide_master] if show_master else []) + [layout]
        for host in hosts:
            def _hairline(sh):
                try:
                    return (sh.shape_type == 13 and sh.height and sh.height * _sy(prs) <= 12
                            and (sh.left < 0 or (sh.left + (sh.width or 0)) * _sx(prs) > CANVAS_W + 4))
                except Exception:
                    return False
            deco = [sh for sh in host.shapes
                    if not getattr(sh, "is_placeholder", False) and not _hairline(sh)]
            if deco:
                parts.append(_walk_shapes(deco, tf, sx, theme_map, font_scheme, recolor, flip,
                                          prs, None, top=False, flip_target=(target if flip else None)))
    except Exception:
        pass
    parts.append(_walk_shapes(slide.shapes, tf, sx, theme_map, font_scheme, recolor, flip, prs, raster_shapes, suppress_names=suppress_names, flip_target=(target if flip else None), logo_swaps=logo_swaps, dark_backers=dark_backers, dark_display=(target == "dark")))
    return "\n".join(p for p in parts if p)
