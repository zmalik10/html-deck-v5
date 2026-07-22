"""
Ingest past decks into the reference library's deck archive (and promote the
recognizable About/bio slides into the curated catalog).

Scans a folder for .pptx / .pdf / .html decks, extracts per-slide
title / text / role / entities into layouts/reference-library/archive.json, and
promotes recognizable About/bio slides into catalog.json with a computed
content_hash. This is how the "reinvent nothing already published" rule gets its
raw material: PLAN then searches the archive + catalog via reference_search.py.

Idempotent: each deck stores a content_hash over its ordered per-slide hashes.
Re-ingesting an unchanged deck is a no-op; a changed deck replaces its own slides
only. --dry-run reports the plan without writing.

Usage:
    python engine/ingest_decks.py                      # default --src (below)
    python engine/ingest_decks.py --src "<folder>"     # scan a specific folder
    python engine/ingest_decks.py --dry-run            # report, write nothing
    python engine/ingest_decks.py --no-promote         # archive only, skip catalog promotion

Optional parsers (best-effort, graceful if missing):
    python-pptx (pptx)  ·  pypdf / PyPDF2 / pdfminer.six (pdf)
    .pptx also falls back to raw XML text extraction with only the stdlib.
"""
import argparse, hashlib, json, os, re, sys, zipfile
from datetime import datetime, timezone

HERE = os.path.dirname(os.path.abspath(__file__))
SKILL = os.path.dirname(HERE)
REFDIR = os.path.join(SKILL, "layouts", "reference-library")
CATALOG = os.path.join(REFDIR, "catalog.json")
ARCHIVE = os.path.join(REFDIR, "archive.json")

DEFAULT_SRC = r"C:\Users\rowan\OneDrive - SMARTBUILD Construction Solutions\Documents\1. Personal\Files"

KNOWN_ENTITIES = [
    "SmartBuild", "SMARTBUILD", "smrt-E", "smrtAEC", "smrt-AE", "smrt-GC", "smrt-SUB",
    "Rowan Steel Hall", "Rowan", "Zulq", "Vinay",
]


# --- content_hash: MUST stay identical to engine/reference_search.py ----------
def content_hash(parts):
    """sha256:16 over an ordered, stripped list of strings."""
    norm = [(p or "").strip() for p in parts]
    payload = json.dumps(norm, ensure_ascii=False, separators=(",", ":"))
    return "sha256:" + hashlib.sha256(payload.encode("utf-8")).hexdigest()[:16]


def _now():
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat()


def _slug(s):
    return re.sub(r"[^a-z0-9]+", "-", (s or "").lower()).strip("-")[:48] or "deck"


def _uuid_from(seed):
    """Deterministic, schema-valid UUID from a seed string. Used only as a placeholder
    slide_uuid in the catalog — PLAN re-mints every uuid on reuse, so the value is a
    stable filler, not an identity."""
    h = hashlib.sha256((seed or "x").encode("utf-8")).hexdigest()
    return "%s-%s-4%s-8%s-%s" % (h[:8], h[8:12], h[13:16], h[17:20], h[20:32])


def _load(path, default):
    if not os.path.exists(path):
        return default
    with open(path, encoding="utf-8") as f:
        return json.load(f)


def _atomic_write(path, data):
    tmp = path + ".tmp"
    with open(tmp, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)
    os.replace(tmp, path)


# --- extraction: one function per format, each -> [{title, text}] -------------
def _split_title_body(lines):
    lines = [l.strip() for l in lines if l and l.strip()]
    if not lines:
        return "", ""
    title = lines[0][:120]
    body = " ".join(lines[1:]).strip()
    return title, body[:1200]


def extract_pptx(path):
    slides = []
    # Preferred: python-pptx (keeps per-shape text).
    try:
        from pptx import Presentation
        prs = Presentation(path)
        for slide in prs.slides:
            lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(r.text for r in para.runs).strip()
                        if t:
                            lines.append(t)
            title, body = _split_title_body(lines)
            slides.append({"title": title, "text": body})
        return slides
    except ImportError:
        pass
    except Exception as e:
        print("  ! python-pptx failed (%s); trying raw XML" % e)
    # Fallback: unzip and read <a:t> runs from ppt/slides/slideN.xml.
    try:
        with zipfile.ZipFile(path) as z:
            names = [n for n in z.namelist()
                     if re.match(r"ppt/slides/slide\d+\.xml$", n)]
            names.sort(key=lambda n: int(re.search(r"slide(\d+)\.xml", n).group(1)))
            for n in names:
                xml = z.read(n).decode("utf-8", "ignore")
                runs = re.findall(r"<a:t>(.*?)</a:t>", xml, re.S)
                lines = [re.sub(r"<[^>]+>", "", r).strip() for r in runs]
                title, body = _split_title_body(lines)
                slides.append({"title": title, "text": body})
    except Exception as e:
        print("  ! could not read pptx %s: %s" % (os.path.basename(path), e))
    return slides


def extract_pdf(path):
    slides = []
    try:
        from pypdf import PdfReader
    except ImportError:
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            PdfReader = None
    if PdfReader:
        try:
            for page in PdfReader(path).pages:
                title, body = _split_title_body((page.extract_text() or "").splitlines())
                slides.append({"title": title, "text": body})
            return slides
        except Exception as e:
            print("  ! pypdf failed on %s: %s" % (os.path.basename(path), e))
    try:
        from pdfminer.high_level import extract_text
        text = extract_text(path) or ""
        for chunk in text.split("\x0c"):  # form-feed = page break
            if chunk.strip():
                title, body = _split_title_body(chunk.splitlines())
                slides.append({"title": title, "text": body})
        return slides
    except ImportError:
        print("  ! no PDF parser (pip install pypdf) — skipping %s" % os.path.basename(path))
    except Exception as e:
        print("  ! pdfminer failed on %s: %s" % (os.path.basename(path), e))
    return slides


def _strip_tags(html):
    html = re.sub(r"<(script|style)[^>]*>.*?</\1>", " ", html, flags=re.S | re.I)
    html = re.sub(r"<br\s*/?>", "\n", html, flags=re.I)
    html = re.sub(r"</(p|div|h[1-6]|li|section)>", "\n", html, flags=re.I)
    text = re.sub(r"<[^>]+>", " ", html)
    text = re.sub(r"&nbsp;", " ", text)
    text = re.sub(r"&amp;", "&", text)
    return text


def extract_html(path):
    with open(path, encoding="utf-8", errors="ignore") as f:
        html = f.read()
    slides = []
    sections = re.findall(r"<section\b[^>]*>(.*?)</section>", html, re.S | re.I)
    if sections:
        heads = re.findall(r"<section\b([^>]*)>", html, re.I)
        for attrs, inner in zip(heads, sections):
            topic = ""
            m = re.search(r'data-topic="([^"]*)"', attrs)
            if m:
                topic = m.group(1)
            lines = [l for l in _strip_tags(inner).splitlines()]
            title, body = _split_title_body(lines)
            slides.append({"title": topic or title, "text": body})
        return slides
    # No slide sections — treat the doc as a single "slide".
    title = ""
    m = re.search(r"<title>(.*?)</title>", html, re.S | re.I)
    if m:
        title = re.sub(r"\s+", " ", m.group(1)).strip()
    _t, body = _split_title_body(_strip_tags(html).splitlines())
    slides.append({"title": title or _t, "text": body})
    return slides


EXTRACTORS = {".pptx": extract_pptx, ".pdf": extract_pdf, ".html": extract_html, ".htm": extract_html}


# --- classification -----------------------------------------------------------
def classify_role(index, total, title, text):
    t = (title + " " + text).lower()
    ttl = title.lower()
    if re.search(r"\babout\b|who we are|our (company|story|mission)", t):
        return "about"
    if re.search(r"\bbio\b|meet the team|leadership|founder|chief |,\s*(ceo|coo|cto|cfo)\b|\b(ceo|coo|cto|cfo)\b", t):
        return "bio"
    if index == 0:
        return "cover"
    if re.search(r"thank you|next steps|let'?s talk|contact us|get in touch", t):
        return "closing"
    if re.search(r"our mission|mission statement", t):
        return "mission"
    if len(ttl) <= 60 and len(text) < 40:
        return "section"
    return "content"


def extract_entities(title, text):
    hay = title + " " + text
    found = []
    for e in KNOWN_ENTITIES:
        if re.search(r"\b" + re.escape(e) + r"\b", hay, re.I) and e not in found:
            found.append(e)
    m = re.search(r"about\s+([A-Z][\w&.-]+(?:\s+[A-Z][\w&.-]+)?)", hay)
    if m and m.group(1) not in found:
        found.append(m.group(1))
    return found


# --- promotion: build a catalog fragment from an archived About/bio slide -----
def _fragment_from_slide(slide, role):
    family = {"about": "NM-01", "bio": "WT-05"}.get(role, "custom")
    blocks = [{"block_uuid": "00000000-0000-4000-8000-000000000001",
               "type": "headline", "text": slide["title"] or "(untitled)"}]
    if slide.get("text"):
        blocks.append({"block_uuid": "00000000-0000-4000-8000-000000000002",
                       "type": "body", "text": slide["text"]})
    return {
        "topic": slide["title"] or role,
        "layout": {"family": family, "shape_tags": [role]},
        "group": role,
        "continues": False,
        "content_blocks": blocks,
        "icon_intent": [],
    }


def entry_content_hash(entry):
    """Same recipe reference_search uses: active block texts + fact values/texts."""
    frag = entry.get("planFragment", {}) or {}
    parts = [b.get("text", "") for b in frag.get("content_blocks", [])
             if b.get("status", "active") != "deleted"]
    for f in entry.get("approved_facts", []):
        parts.append(f.get("value") or f.get("text", ""))
    return content_hash(parts)


def _match_canonical(cat, role, entities):
    ents = [e.lower() for e in entities]
    for e in cat.get("entries", []):
        if not e.get("canonical") or e.get("role") != role:
            continue
        cat_ents = [x.lower() for x in e.get("entities", [])]
        if any(a in cat_ents or any(a in c or c in a for c in cat_ents) for a in ents):
            return e
    return None


def promote(cat, arc_slide, role, entities, deck_title, source_path, dry):
    """Promote an About/bio archive slide into the catalog. Returns a change note or None."""
    frag = _fragment_from_slide(arc_slide, role)
    target = _match_canonical(cat, role, entities)

    if target:
        candidate = dict(target)
        sid = (target.get("planFragment") or {}).get("slide_uuid") or _uuid_from(target.get("ref_id", ""))
        candidate["planFragment"] = {**frag, "slide_uuid": sid}
        new_hash = entry_content_hash(candidate)
        if target.get("content_hash") == new_hash and target.get("status") == "published":
            return None  # already up to date — idempotent
        note = "update %s (draft->published)" % target["ref_id"] if target.get("status") == "draft" \
               else "refresh %s" % target["ref_id"]
        if not dry:
            target["planFragment"] = candidate["planFragment"]
            target["status"] = "published"
            target["content_hash"] = new_hash
            target["provenance"] = {
                "origin": "ingest", "authored_by": None,
                "source_deck": deck_title, "source_path": source_path,
                "ingested_at": _now(),
                "note": "Promoted from ingested published slide; replaced the draft seed.",
            }
        return note

    # No canonical target — add a non-canonical catalog candidate.
    ref_id = "REF-ingested-%s-%s" % (role, _slug(arc_slide["title"] or deck_title))
    if any(e.get("ref_id") == ref_id for e in cat.get("entries", [])):
        return None
    entry = {
        "ref_id": ref_id, "title": arc_slide["title"] or ("%s (%s)" % (role, deck_title)),
        "role": role, "entities": entities, "topics": [], "audience": ["internal"],
        "canonical": False, "status": "published", "layout_hint": frag["layout"]["family"],
        "content_hash": "sha256:PENDING", "approved_facts": [],
        "provenance": {"origin": "ingest", "authored_by": None, "source_deck": deck_title,
                       "source_path": source_path, "ingested_at": _now(),
                       "note": "Auto-added from ingested deck; not canonical."},
        "used_in": [],
        "planFragment": {**frag, "slide_uuid": "00000000-0000-4000-8000-%012d" % (len(cat.get("entries", [])) + 1)},
    }
    entry["content_hash"] = entry_content_hash(entry)
    if not dry:
        cat.setdefault("entries", []).append(entry)
    return "add %s" % ref_id


def extract_images(pptx_path, coll_name, dry):
    """RC10: pull every picture out of a collection PPTX into assets/images/ (deduped by
    content hash, named by slide+index) and register each as an OWNED image asset in
    libraries/images/catalog.json with provenance. This gives the product template real
    product screenshots to draw from instead of a CSS mockup. Idempotent by content hash."""
    try:
        from pptx import Presentation
    except Exception:
        print("  ! python-pptx required for --extract-images"); return []
    imgs_dir = os.path.join(SKILL, "assets", "images", "extracted")
    cat_path = os.path.join(SKILL, "libraries", "images", "catalog.json")
    cat = _load(cat_path, {"owned": {"entries": []}})
    owned = cat.setdefault("owned", {}).setdefault("entries", [])
    by_hash = {e.get("blob_sha"): e for e in owned if e.get("blob_sha")}
    prs = Presentation(pptx_path)
    notes, seen = [], set()
    if not dry:
        os.makedirs(imgs_dir, exist_ok=True)

    def _walk(shapes, si):
        for i, sh in enumerate(shapes):
            if sh.shape_type == 6:
                _walk(sh.shapes, si); continue
            if sh.shape_type != 13:
                continue
            try:
                blob = sh.image.blob; ext = sh.image.ext or "png"
            except Exception:
                continue
            h = hashlib.sha256(blob).hexdigest()[:16]
            if h in seen:
                continue
            seen.add(h)
            if h in by_hash:
                continue                                    # already registered — idempotent
            try:
                from PIL import Image; import io
                im = Image.open(io.BytesIO(blob)); w, hh = im.size
            except Exception:
                w = hh = 0
            fn = "%s-s%02d-i%02d.%s" % (_slug(coll_name), si + 1, i, ext)
            rel = "assets/images/extracted/%s" % fn
            if not dry:
                with open(os.path.join(imgs_dir, fn), "wb") as f:
                    f.write(blob)
            entry = {
                "id": "img-%s-s%02d-%02d" % (_slug(coll_name), si + 1, i),
                "tags": ["extracted", "%s-s%02d" % (_slug(coll_name), si + 1)],
                "mood": "product screenshot", "file": rel, "blob_sha": h,
                "dims": [w, hh], "status": "seeded",
                "provenance": {"kind": "image", "source": "SmartBuild collection: %s" % coll_name,
                               "author": "SmartBuild", "license": "owned",
                               "requires_attribution": False, "approved_for_client": True},
            }
            owned.append(entry)
            notes.append("    -> image: %s (%dx%d) from slide %d" % (rel, w, hh, si + 1))

    for si, slide in enumerate(prs.slides):
        _walk(slide.shapes, si)
    if not dry:
        _atomic_write(cat_path, cat)
    return notes


def _load_curation():
    """Curation overlay: human decisions (canonical, corrected role, entities, tags) that
    must SURVIVE re-ingest. Keyed by a case-insensitive title substring so it holds even
    when the collection is reordered. See layouts/reference-library/curation.json."""
    return _load(os.path.join(REFDIR, "curation.json"), {"rules": []}).get("rules", [])


def _apply_curation(entry, rules):
    title = (entry.get("title") or "").lower()
    for r in rules:
        mt = (r.get("match_title") or "").lower()
        if mt and mt in title:
            entry.update(r.get("set", {}))
    return entry


def _curation():
    return _load(os.path.join(REFDIR, "curation.json"), {})


def _suppressions(index1):
    """Curation-driven per-slide shape suppression (review pin-2): curation.json may carry
    shape_suppressions: [{"slide": <1-based index>, "names": [shape names]}]."""
    for r in _curation().get("shape_suppressions", []):
        if r.get("slide") == index1:
            return set(r.get("names", []))
    return None


def _logo_files(index1):
    """Curation-driven per-slide brand-logo swaps (W1a): returns {shape_name: logo_filename}
    from curation.json logo_swaps. ingest builds the theme-correct SVG via render_reference."""
    for r in _curation().get("logo_swaps", []):
        if r.get("slide") == index1:
            return dict(r.get("shapes", {}))
    return None


def _image_files(index1):
    """Curation-driven per-slide RASTER swaps (cycle-10 pin-5): returns {shape_name:
    asset_filename} from curation.json image_swaps. ingest builds the <img> via
    render_reference.img_swap. Unlike logo_swaps these are theme-independent PNGs."""
    for r in _curation().get("image_swaps", []):
        if r.get("slide") == index1:
            return dict(r.get("shapes", {}))
    return None


def _dark_backers(index1):
    """Curation-driven per-slide dark-mode logo backers (cycle-10 pin-6): shape names that
    get a soft light chip drawn behind them when the displayed theme is dark."""
    for r in _curation().get("dark_backers", []):
        if r.get("slide") == index1:
            return set(r.get("names", []))
    return None


def _fragment_rules(index1):
    """Curation-driven per-slide rebuild/append fragments (W1a): native HTML overrides
    (s04 rebuild, s05 chart, s11 dark overlay) that survive re-ingest. See curation.json
    fragments: [{slide, mode: replace|append|append-dark, file, clip_dark?}]."""
    return [r for r in _curation().get("fragments", []) if r.get("slide") == index1]


def _apply_fragments(frag_html, theme, rules, fragdir):
    """Apply curation fragment rules to a rendered slide fragment for one theme, matching
    build_ref_review.py exactly. replace=swap whole slide; append=add to both themes;
    append-dark=on dark only, clip the matched raster then append the overlay."""
    for r in rules:
        mode = r.get("mode", "append")
        path = os.path.join(fragdir, r.get("file", ""))
        try:
            content = open(path, encoding="utf-8").read()
        except Exception as e:
            print("    ! fragment %s missing (%s)" % (r.get("file"), e)); continue
        if mode == "replace":
            frag_html = content
        elif mode == "append":
            frag_html = frag_html + content
        elif mode == "append-dark" and theme == "dark":
            clip = r.get("clip_dark")
            if clip:
                def _clip(m):
                    return m.group(0).replace('style="', 'style="' + clip["insert"], 1)
                frag_html = re.sub(clip["match"], _clip, frag_html, count=1)
            frag_html = frag_html + content
    return frag_html


def _render_verbatim(entry, pptx_path, index1, refdir, dry, suppress_names=None,
                     logo_files=None, fragment_rules=None, image_files=None, dark_backers=None):
    """Pre-render a canonical entry to LOCKED verbatim HTML in both themes (native +
    generated counterpart) so build.py can drop it in word/image/style-for-style and
    only swap light<->dark. Applies the SAME curation overlays (shape suppressions, brand
    logo swaps, native rebuild/append fragments) as the exec-approved review deck so catalog
    renders match it exactly. Best-effort: leaves the entry as text-reuse if rendering fails."""
    try:
        import render_reference as RR
        from pptx import Presentation
        idx0 = index1 - 1
        prs = Presentation(pptx_path)
        tm = RR.load_theme_map(pptx_path)
        fs = RR.load_font_scheme(pptx_path)
        native = RR.native_theme_of(prs, prs.slides[idx0], tm)
        # RC3: COM-rasterize shapes that can't be faithfully rebuilt (groups/diagrams/
        # SmartArt/charts) incl. unrenderable pictures (linked/svg-graphic parts, pin-12).
        rasters = RR.rasterize_shapes(pptx_path, idx0, extra_idx=RR.unrenderable_pics(prs, idx0))
        if rasters:
            print("    [fidelity] slide %d: %d complex shape(s) rasterized via COM" % (index1, len(rasters)))
        fragdir = os.path.join(refdir, "fragments")
        fragment_rules = fragment_rules or []
        rid = entry["ref_id"]
        rendered = os.path.join(refdir, "rendered")
        rels, modes = {}, {}
        for theme in ("light", "dark"):
            swaps = {}
            if logo_files:
                swaps.update({n: RR.logo_svg(f, dark=(theme == "dark")) for n, f in logo_files.items()})
            if image_files:
                swaps.update({n: RR.img_swap(f) for n, f in image_files.items()})
            swaps = swaps or None
            frag = RR.render_slide(prs, idx0, tm, fs, theme=theme, raster_shapes=rasters,
                                   suppress_names=suppress_names, logo_swaps=swaps,
                                   dark_backers=dark_backers)
            frag = _apply_fragments(frag, theme, fragment_rules, fragdir)
            modes[theme] = "faithful" if theme == native else (
                "image-swap" if any(sh.shape_type == 13 and RR._is_fullbleed(sh, prs) for sh in prs.slides[idx0].shapes)
                else "recolor")
            rel = os.path.join("rendered", "%s.%s.html" % (rid, theme))
            if not dry:
                os.makedirs(rendered, exist_ok=True)
                with open(os.path.join(refdir, rel), "w", encoding="utf-8") as f:
                    f.write(frag)
            rels[theme] = rel.replace("\\", "/")
        entry.update({
            "render_mode": "verbatim-locked", "locked": True, "exec_approved": True,
            "native_theme": native, "render": rels,
            "render_modes": modes,                        # RC1: which transform each theme used
            "rasterized_shapes": sorted(rasters.keys()),  # RC3: shapes kept as raster (never omitted)
            "theme_note": "STYLE is locked (word/image/style-for-style); colours are converted to the "
                          "brand guide for the counterpart theme (RC1). Complex diagrams are COM-rasterized (RC3).",
        })
        return True
    except Exception as e:
        print("    ! verbatim render failed for %s: %s" % (entry.get("ref_id"), e))
        return False


def ingest_collection(cat, slides, coll_name, source_path, dry, rules=None, pptx_path=None):
    """Every slide in a curated collection -> a published catalog entry. Idempotent by
    replacement: drops this collection's prior entries, then re-adds all current slides.
    Applies the curation overlay so canonical/role decisions survive re-ingest. Canonical
    entries are additionally pre-rendered to LOCKED verbatim HTML (both themes)."""
    rules = rules or []
    if not dry:
        cat["entries"] = [e for e in cat.get("entries", [])
                          if (e.get("provenance") or {}).get("source_collection") != coll_name]
    notes = []
    for i, s in enumerate(slides, 1):
        role = classify_role(i - 1, len(slides), s["title"], s["text"])
        ents = extract_entities(s["title"], s["text"]) or ["SmartBuild"]
        frag = _fragment_from_slide(s, role)
        frag["slide_uuid"] = _uuid_from("%s#%d" % (coll_name, i))
        entry = {
            "ref_id": "REF-%s-s%02d" % (_slug(coll_name), i),
            "title": s["title"] or ("%s slide %d" % (coll_name, i)),
            "role": role, "entities": ents, "topics": [], "audience": ["internal", "client"],
            "canonical": False, "status": "published", "layout_hint": frag["layout"]["family"],
            "content_hash": "sha256:PENDING", "approved_facts": [],
            "provenance": {"origin": "collection", "source_collection": coll_name,
                           "source_index": i, "source_path": source_path, "ingested_at": _now(),
                           "note": "Curated CURRENT reference slide — every slide in this collection is a keeper. "
                                   "Mark canonical=true on the always-reuse ones (About/bios) during curation."},
            "used_in": [], "planFragment": frag,
        }
        _apply_curation(entry, rules)                       # human overrides survive re-ingest
        if pptx_path:   # pin-25: EVERY collection slide gets locked light+dark renders
            _render_verbatim(entry, pptx_path, i, REFDIR, dry, suppress_names=_suppressions(i),
                             logo_files=_logo_files(i), fragment_rules=_fragment_rules(i),
                             image_files=_image_files(i), dark_backers=_dark_backers(i))
        entry["content_hash"] = entry_content_hash(entry)
        if not dry:
            cat.setdefault("entries", []).append(entry)
        flag = (" *CANON,LOCKED*" if entry.get("render_mode") == "verbatim-locked"
                else " *CANON*" if entry.get("canonical") else "")
        notes.append("    -> catalog: %s  [%s]%s  %s" % (entry["ref_id"], entry.get("role"), flag, (s["title"] or "")[:44]))
    return notes


# --- main ---------------------------------------------------------------------
def find_decks(src):
    hits = []
    for root, _dirs, files in os.walk(src):
        for fn in files:
            if fn.startswith("~$") or fn.startswith("."):
                continue
            if os.path.splitext(fn)[1].lower() in EXTRACTORS:
                hits.append(os.path.join(root, fn))
    return sorted(hits)


def main():
    ap = argparse.ArgumentParser(description="Ingest past decks into the reference-library archive + catalog.")
    ap.add_argument("--src", default=DEFAULT_SRC, help="folder to scan for .pptx/.pdf/.html decks")
    ap.add_argument("--dry-run", action="store_true", help="report the plan; write nothing")
    ap.add_argument("--no-promote", action="store_true", help="archive only; skip catalog promotion")
    ap.add_argument("--extract-images", action="store_true",
                    help="RC10: extract every picture from each collection PPTX into assets/images/ "
                         "and register them as owned image assets (for the product template's data-image slot).")
    ap.add_argument("--collection", action="store_true",
                    help="treat each source file as a CURATED COLLECTION of current, reusable slides: "
                         "every slide becomes a published catalog entry (not archived-only). "
                         "Re-ingest replaces that collection's entries (idempotent by source_collection).")
    args = ap.parse_args()

    if not os.path.isdir(args.src):
        print("Source folder not found:", args.src); sys.exit(1)

    arc = _load(ARCHIVE, {"summary": {}, "decks": [], "slides": []})
    cat = _load(CATALOG, {"summary": {}, "entries": []})
    decks_by_path = {d["source_path"]: d for d in arc.get("decks", [])}

    files = find_decks(args.src)
    print("Scanning %s\nFound %d candidate file(s).\n" % (args.src, len(files)))

    added_decks = changed = skipped = 0
    promo_notes = []

    for path in files:
        ext = os.path.splitext(path)[1].lower()
        rel = os.path.basename(path)
        slides = EXTRACTORS[ext](path)
        if not slides:
            print("  - %s: no slides extracted, skipping" % rel); continue

        if args.collection:
            coll_name = os.path.splitext(rel)[0]
            print("  * %s: %d slides -> CURATED COLLECTION '%s' (every slide = published catalog entry)"
                  % (rel, len(slides), coll_name))
            promo_notes.extend(ingest_collection(cat, slides, coll_name, path, args.dry_run, _load_curation(), pptx_path=path))
            if args.extract_images and ext == ".pptx":
                promo_notes.extend(extract_images(path, coll_name, args.dry_run))
            changed += 1
            continue

        total = len(slides)
        slide_hashes = [content_hash([s["title"], s["text"]]) for s in slides]
        deck_hash = content_hash(slide_hashes)

        prev = decks_by_path.get(path)
        if prev and prev.get("content_hash") == deck_hash:
            skipped += 1
            print("  = %s: unchanged (%d slides), skip" % (rel, total)); continue

        deck_id = "%s-%s" % (_slug(os.path.splitext(rel)[0]),
                             hashlib.sha256(path.encode("utf-8")).hexdigest()[:8])
        deck_title = os.path.splitext(rel)[0]
        verb = "changed" if prev else "new"
        print("  %s %s: %d slides [%s]" % ("~" if prev else "+", rel, total, verb))

        # Build slide records.
        new_slides = []
        for i, (s, sh) in enumerate(zip(slides, slide_hashes)):
            role = classify_role(i, total, s["title"], s["text"])
            ents = extract_entities(s["title"], s["text"])
            rec = {"archive_id": "%s#%d" % (deck_id, i + 1), "deck_id": deck_id,
                   "index": i + 1, "title": s["title"], "text": s["text"],
                   "role": role, "entities": ents, "content_hash": sh, "used_in": []}
            # preserve any prior reuse stamps for a stable archive_id
            if prev:
                old = next((x for x in arc["slides"]
                            if x.get("archive_id") == rec["archive_id"]), None)
                if old and old.get("used_in"):
                    rec["used_in"] = old["used_in"]
            new_slides.append(rec)

            if role in ("about", "bio") and not args.no_promote:
                note = promote(cat, s, role, ents or ["SmartBuild"], deck_title, path, args.dry_run)
                if note:
                    promo_notes.append("    -> catalog: %s (from %s s%d)" % (note, rel, i + 1))

        if not args.dry_run:
            arc["slides"] = [x for x in arc["slides"] if x.get("deck_id") != deck_id] + new_slides
            deck_rec = {"deck_id": deck_id, "source_path": path,
                        "source_type": ext.lstrip("."), "title": deck_title,
                        "slide_count": total, "content_hash": deck_hash, "ingested_at": _now()}
            arc["decks"] = [d for d in arc["decks"] if d.get("source_path") != path] + [deck_rec]
        added_decks += 0 if prev else 1
        changed += 1 if prev else 0

    for n in promo_notes:
        print(n)

    if not args.dry_run:
        arc["summary"].update({
            "dataset_version": arc["summary"].get("dataset_version", "ref-archive-v1"),
            "kind": "deck-archive",
            "deck_count": len(arc["decks"]), "slide_count": len(arc["slides"]),
            "last_ingest": _now(), "last_src": args.src,
        })
        cat["summary"]["entry_count"] = len(cat.get("entries", []))
        cat["summary"]["canonical_count"] = sum(1 for e in cat["entries"] if e.get("canonical"))
        _atomic_write(ARCHIVE, arc)
        _atomic_write(CATALOG, cat)

    print("\n%s: %d new deck(s), %d changed, %d unchanged. Archive now %d deck(s) / %d slide(s)."
          % ("DRY-RUN (no writes)" if args.dry_run else "Wrote",
             added_decks, changed, skipped,
             len(arc.get("decks", [])), len(arc.get("slides", []))))


if __name__ == "__main__":
    main()
