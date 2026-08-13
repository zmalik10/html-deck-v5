"""
SmartBuild Deck v5 — export_pptx.py  (the OFFICIAL engine-side PPTX EXPORT)

Replaces the browser (deck.js) PPTX path. Two slide classes, two strategies:

  * REFERENCE slides (plan entry has ``reference.ref_id``): the perfect native
    representation already exists — the ORIGINAL slide in the source ``.pptx``. We
    COPY it (PowerPoint COM ``InsertFromFile``) into the output deck at the right
    position, so every shape/group/font/image stays native and movable. When the deck
    theme differs from the entry's ``native_theme`` we run a BRAND RECOLOR pass over the
    copied slide's XML, mapping every colour through the SAME table the HTML renderer
    uses — ``render_reference._map_color`` / ``_map_text_color`` (imported, never
    duplicated). Curated overrides (shape suppressions, brand-logo swaps) are re-applied
    natively so the PPTX matches the exec-approved HTML render.

  * AUTHORED slides: reconstructed as native python-pptx objects from a DOM manifest
    (``window.__exportManifest()`` in deck.js, pulled via headless Chrome/Playwright):
    shapes, images, svgs, text (with per-run styles + z-order), and native charts.
    Text whose container is small (<=72px) is attached to its shape (text-in-shape) so
    numbers stay centred in circles. See ``build_authored_slide``.

Design goals: engine-side (testable, no CDN roulette), deterministic, and readback-
verifiable — the fidelity harness asserts >3 native shapes per reference slide, every
visible string exactly once as an editable run, no slide-covering rasters, drift <=1%.

Usage:
    python export_pptx.py --skill-path <dir> --plan plan.json --slides-html <built presentation.html> \
        --out deck.pptx [--theme dark|light] [--no-authored] [--refs-only]
"""
import argparse, json, os, sys, time, re

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import render_reference as RR      # the ONE colour-mapping owner (imported, not copied)

EMU = 914400
LOGW, LOGH = 1280.0, 720.0
PPTW_IN, PPTH_IN = 13.333, 7.5
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS = {"a": A_NS}


def load_json(p):
    with open(p, encoding="utf-8") as f:
        return json.load(f)


# ----------------------------------------------------------------- catalog lookup
def catalog_entries(sp):
    catp = os.path.join(sp, "layouts", "reference-library", "catalog.json")
    if not os.path.exists(catp):
        return {}
    return {e["ref_id"]: e for e in load_json(catp).get("entries", [])}


def curation(sp):
    cp = os.path.join(sp, "layouts", "reference-library", "curation.json")
    return load_json(cp) if os.path.exists(cp) else {}


# =====================================================================================
#  R1 — BRAND-LOGO SWAPS in copied reference slides
#  The curated logo_swaps replace raster brand marks in the exec-approved review deck
#  with crisp vector logos. Rowan pinned this repeatedly: the marks must be the real
#  brand SVGs, not the source deck's low-res rasters. Here we render each brand SVG to a
#  high-resolution transparent PNG (>=3x the shape box, Playwright screenshots the inline
#  SVG produced by the SHARED render_reference.logo_svg — same dark grey-lift as the HTML
#  catalog render) and drop it into the copied slide at the SAME box. Cached on disk.
# =====================================================================================
class LogoRasterizer:
    """Render assets/logos/*.svg to transparent PNG bytes at a requested pixel size, via
    a single headless Chromium instance (opened lazily, reused for every logo). Disk-cached
    under <skill>/.local/logo-cache so repeat exports don't re-rasterize."""

    def __init__(self, sp):
        self.sp = sp
        self.cache_dir = os.path.join(sp, ".local", "logo-cache")
        os.makedirs(self.cache_dir, exist_ok=True)
        self._pw = self._browser = self._page = None

    def _ensure(self):
        if self._page is not None:
            return
        from playwright.sync_api import sync_playwright
        self._pw = sync_playwright().start()
        self._browser = self._pw.chromium.launch()
        self._page = self._browser.new_page(device_scale_factor=1)

    def png(self, svg_file, dark, w_px, h_px):
        w_px = max(2, int(round(w_px)))
        h_px = max(2, int(round(h_px)))
        slug = re.sub(r"[^a-z0-9]+", "", os.path.basename(svg_file).lower())
        key = "%s-%s-%dx%d.png" % (slug, "d" if dark else "l", w_px, h_px)
        fp = os.path.join(self.cache_dir, key)
        if os.path.exists(fp):
            with open(fp, "rb") as f:
                return f.read()
        self._ensure()
        svg = RR.logo_svg(svg_file, dark=dark)   # SHARED helper (id de-collision + dark grey-lift)
        html = (
            "<!doctype html><html><head><meta charset='utf-8'><style>"
            "*{margin:0;padding:0}html,body{background:transparent}"
            "#b{width:%dpx;height:%dpx}#b svg{display:block;width:100%%;height:100%%}"
            "</style></head><body><div id='b'>%s</div></body></html>" % (w_px, h_px, svg))
        self._page.set_viewport_size({"width": w_px, "height": h_px})
        self._page.set_content(html, wait_until="networkidle")
        el = self._page.query_selector("#b")
        raw = el.screenshot(omit_background=True, type="png")
        with open(fp, "wb") as f:
            f.write(raw)
        return raw

    def close(self):
        try:
            if self._browser is not None:
                self._browser.close()
            if self._pw is not None:
                self._pw.stop()
        except Exception:
            pass


def apply_logo_swaps(slide, shapes_map, theme, rast):
    """Replace named raster picture shapes with high-res brand-SVG PNGs at the same box.
    dark theme -> dark grey-lifted variant. Returns count swapped."""
    import io
    dark = (theme == "dark")
    n = 0
    for sh in list(slide.shapes):
        if sh.name in shapes_map and sh.left is not None and sh.width and sh.height:
            L, T, W, H = sh.left, sh.top, sh.width, sh.height
            w_px = W / EMU * 96 * 3
            h_px = H / EMU * 96 * 3
            png = rast.png(shapes_map[sh.name], dark, w_px, h_px)
            sh._element.getparent().remove(sh._element)     # drop the raster mark
            slide.shapes.add_picture(io.BytesIO(png), L, T, W, H)
            n += 1
    return n


# =====================================================================================
#  R2 — NATIVE PPTX BUILDERS for curated-fragment reference slides (s04 rebuild, s05 chart)
#  When a deck reuses these reference slides, we must NOT ship the raster chart or a
#  picture-of-text. Instead we render the curated fragment (the exec-approved native
#  rebuild) in headless Chromium, extract a native-object manifest (panels/bar rects as
#  shapes, all text as editable runs, the SVG wordmark rasterized), and build the slide
#  with python-pptx via build_authored_slide — the same path authored slides use.
# =====================================================================================
_FRAG_EXTRACT_JS = r"""() => {
  const stage = document.getElementById('stage');
  const sr = stage.getBoundingClientRect(), sc = sr.width / 1280;
  const H = (c) => { const m = (c||'').match(/rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([\d.]+))?/);
    if (!m) return null; const a = m[4]===undefined?1:parseFloat(m[4]);
    return { hex: [1,2,3].map(i=>(+m[i]).toString(16).padStart(2,'0')).join(''), alpha: a }; };
  const box = (el) => { const r = el.getBoundingClientRect();
    return { x:(r.left-sr.left)/sc, y:(r.top-sr.top)/sc, w:r.width/sc, h:r.height/sc }; };
  const INLINE = ['SPAN','B','STRONG','EM','I','BR'];
  const isTextUnit = (el) => { if(!(el.innerText||'').trim()) return false;
    for (const k of el.children) if (!INLINE.includes(k.tagName)) return false; return true; };
  const runsOf = (el) => { const out = [];
    (function w(n){ for (const c of n.childNodes){
      if (c.nodeType===3){ let t=c.nodeValue.replace(/\s+/g,' '); if (t.trim()!==''||t===' '){
        const cs=getComputedStyle(n); let tx=t; if(cs.textTransform==='uppercase') tx=tx.toUpperCase();
        out.push({text:tx, color:(H(cs.color)||{hex:'000000'}).hex, bold:parseInt(cs.fontWeight,10)>=600, italic:cs.fontStyle==='italic', underline:(cs.textDecorationLine||cs.textDecoration||'').indexOf('underline')>=0, size:parseFloat(cs.fontSize)}); } }
      else if (c.nodeType===1){ if (c.tagName==='BR'){ if(out.length) out[out.length-1].text+='\n'; } else w(c); } } })(el);
    while (out.length && out[0].text.trim()==='') out.shift();
    while (out.length && out[out.length-1].text.trim()==='') out.pop();
    return out; };
  const alignOf = (cs) => cs.textAlign==='center'?'center':((cs.textAlign==='right'||cs.textAlign==='end')?'right':'left');
  const man = { shapes: [], images: [], charts: [], texts: [], media: [] };
  const sbg = H(getComputedStyle(stage).backgroundColor); if (sbg) man.bg = sbg.hex;
  let z = 0, mi = 0;
  // media (svg/img) -> tagged so the Python side can screenshot each to PNG
  stage.querySelectorAll('svg,img').forEach(el => { const b = box(el); if (b.w<3||b.h<3) return;
    el.setAttribute('data-mi', mi); man.media.push({ mi: mi, box: b }); mi++; });
  // shape + text layers
  (function walk(node){ Array.prototype.forEach.call(node.children, el => {
    const tag = el.tagName.toUpperCase();
    if (tag==='IMG'||tag==='SVG') return;
    const cs = getComputedStyle(el);
    if (cs.display==='none' || cs.visibility==='hidden' || parseFloat(cs.opacity)===0) return;
    const b = box(el);
    let spec = null;
    if (b.w>0.4 && b.h>0.4){ let has=false; spec={box:b, shape:'rect'};
      const bg = H(cs.backgroundColor);
      if (bg && bg.alpha>0.02){ spec.fill=bg.hex; if(bg.alpha<0.99) spec.fillAlpha=bg.alpha; has=true; }
      const bw=parseFloat(cs.borderTopWidth)||0, bc=H(cs.borderTopColor);
      const uniform = bw>0.5 && bc && bc.alpha>0.05
        && cs.borderTopWidth===cs.borderRightWidth && cs.borderTopWidth===cs.borderBottomWidth && cs.borderTopWidth===cs.borderLeftWidth
        && cs.borderTopColor===cs.borderRightColor && cs.borderTopColor===cs.borderBottomColor && cs.borderTopColor===cs.borderLeftColor;
      if (uniform){ spec.line=bc.hex; spec.lineW=bw; has=true; }
      if (has){ const rad=parseFloat(cs.borderTopLeftRadius)||0;
        if (rad>0 && rad>=Math.min(b.w,b.h)/2-1) spec.shape='ellipse'; else if (rad>2) spec.shape='roundRect';
        spec.z = z++; } else spec=null; }
    if (spec) man.shapes.push(spec);
    if (isTextUnit(el)){ const runs=runsOf(el);
      if (runs.length) man.texts.push({ box: box(el), runs: runs, align: alignOf(cs),
        lineHeight: /px$/.test(cs.lineHeight)?parseFloat(cs.lineHeight):null });
      return; }
    walk(el);
  }); })(stage);
  return man;
}"""


def fragment_manifest(fragment_html, theme, bg_hex=None):
    """Render a curated fragment as a standalone 1280x720 stage and return a native-object
    manifest (compatible with build_authored_slide). SVG/img marks are rasterized to PNG
    via element screenshots (python-pptx cannot embed SVG). bg_hex=None -> deck theme bg;
    pass "" (empty) for an overlay fragment that must not paint its own background."""
    from playwright.sync_api import sync_playwright
    if bg_hex is None:
        bg_hex = RR.DARK_BG if theme == "dark" else RR.LIGHT_BG
    stage_bg = ("background:%s;" % bg_hex) if bg_hex else "background:transparent;"
    page_html = (
        "<!doctype html><html data-theme='%s'><head><meta charset='utf-8'><style>"
        "*{margin:0;padding:0;box-sizing:border-box}html,body{width:1280px;height:720px}"
        "#stage{position:relative;width:1280px;height:720px;overflow:hidden;%s"
        "font-family:Montserrat,Arial,sans-serif}</style></head>"
        "<body><div id='stage'>%s</div></body></html>" % (theme, stage_bg, fragment_html))
    import base64
    with sync_playwright() as p:
        b = p.chromium.launch()
        pg = b.new_page(viewport={"width": 1280, "height": 720}, device_scale_factor=1)
        pg.set_content(page_html, wait_until="networkidle")
        man = pg.evaluate(_FRAG_EXTRACT_JS)
        for md in man.get("media", []):
            el = pg.query_selector("[data-mi='%d']" % md["mi"])
            if el is None:
                continue
            try:
                raw = el.screenshot(omit_background=True, type="png")
                man["images"].append({"box": md["box"],
                                      "data": "data:image/png;base64," + base64.b64encode(raw).decode()})
            except Exception:
                pass
        b.close()
    man.pop("media", None)
    if not bg_hex:
        man.pop("bg", None)          # overlay: never repaint the copied slide's background
    return man


def strip_media(slide, min_cov=0.30):
    """Remove large pictures (raster charts / pasted images) and native charts from a
    copied reference slide so the native fragment overlay is the only chart. Used for
    'append'-mode fragment slides (s05). Returns count removed."""
    area = float(SLIDE_W_EMU) * float(SLIDE_H_EMU)
    n = 0
    for sh in list(slide.shapes):
        st = getattr(sh, "shape_type", None)
        remove = False
        if st == 13:                                  # picture
            try:
                if (float(sh.width or 0) * float(sh.height or 0)) / area > min_cov:
                    remove = True
            except Exception:
                pass
        try:
            if getattr(sh, "has_chart", False):        # native chart graphicFrame
                remove = True
        except Exception:
            pass
        if remove:
            sh._element.getparent().remove(sh._element)
            n += 1
    return n


# =====================================================================================
#  COLOUR RECOLOR (brand conversion) — shared table, two consumers
# =====================================================================================
# default DrawingML clrMap (master usually inherits these): text<->dark, bg<->light
_CLRMAP = {"tx1": "dk1", "bg1": "lt1", "tx2": "dk2", "bg2": "lt2",
           "dk1": "dk1", "lt1": "lt1", "dk2": "dk2", "lt2": "lt2",
           "accent1": "accent1", "accent2": "accent2", "accent3": "accent3",
           "accent4": "accent4", "accent5": "accent5", "accent6": "accent6",
           "hlink": "hlink", "folHlink": "folHlink"}


def _theme_scheme(prs_part_theme):
    """{scheme-name: 'RRGGBB'} from a theme XML element."""
    out = {}
    clr = prs_part_theme.find(".//a:clrScheme", NS)
    if clr is None:
        return out
    for child in clr:
        name = child.tag.split("}")[-1]
        srgb = child.find("a:srgbClr", NS)
        sysc = child.find("a:sysClr", NS)
        if srgb is not None:
            out[name] = srgb.get("val")
        elif sysc is not None:
            out[name] = sysc.get("lastClr") or "000000"
    return out


def _scheme_hex(scheme, name):
    """Resolve a schemeClr name (tx1/bg1/accent4/...) to concrete hex via clrMap+scheme."""
    mapped = _CLRMAP.get(name, name)
    # dk1/lt1 in the theme are usually named dk1/lt1 already
    return scheme.get(mapped) or scheme.get(name)


def _is_text_ctx(el):
    a = el.getparent()
    while a is not None:
        tag = a.tag.split("}")[-1]
        if tag in ("rPr", "defRPr", "endParaRPr"):
            return True
        if tag in ("spPr", "bgPr", "grpSpPr", "tcPr", "tblPr"):
            return False
        a = a.getparent()
    return False


def recolor_slide(slide, scheme, target):
    """Map every srgbClr/schemeClr on the slide through the brand table. schemeClr is
    resolved to concrete hex via the source theme, then rewritten as srgbClr so the
    output no longer depends on the (foreign) theme part. Returns count changed."""
    tree = slide.shapes._spTree
    changed = 0
    # 1) resolve schemeClr -> srgbClr (concrete), preserving child modifiers except we
    #    drop the scheme reference; keep alpha/lumMod etc. as children.
    for el in list(tree.findall(".//a:schemeClr", NS)):
        name = el.get("val")
        hexv = _scheme_hex(scheme, name)
        if not hexv:
            continue
        text = _is_text_ctx(el)
        mapped = (RR._map_text_color("#" + hexv, target) if text
                  else RR._map_color("#" + hexv, target)).lstrip("#")
        el.tag = "{%s}srgbClr" % A_NS
        el.set("val", mapped.upper())
        changed += 1
    # 2) map explicit srgbClr
    for el in list(tree.findall(".//a:srgbClr", NS)):
        cur = el.get("val")
        if not cur or len(cur) < 6:
            continue
        text = _is_text_ctx(el)
        mapped = (RR._map_text_color("#" + cur, target) if text
                  else RR._map_color("#" + cur, target)).lstrip("#")
        if mapped.upper() != cur.upper():
            el.set("val", mapped.upper())
            changed += 1
    return changed


def set_background(slide, target):
    """Force the slide background to the brand theme bg (navy on dark, white on light)
    so recoloured content sits on the right surface even if the copied bg was a picture
    or scheme fill we left alone."""
    from pptx.oxml.ns import qn
    hexv = (RR.DARK_BG if target == "dark" else RR.LIGHT_BG).lstrip("#").upper()
    cSld = slide._element.find(qn("p:cSld"))       # <p:bg> lives INSIDE <p:cSld>, not <p:sld>
    # remove existing <p:bg>
    for bg in cSld.findall(qn("p:bg")):
        cSld.remove(bg)
    bg = cSld.makeelement(qn("p:bg"), {})
    bgPr = bg.makeelement(qn("p:bgPr"), {})
    solidFill = bgPr.makeelement(qn("a:solidFill"), {})
    srgb = solidFill.makeelement(qn("a:srgbClr"), {"val": hexv})
    solidFill.append(srgb)
    bgPr.append(solidFill)
    effLst = bgPr.makeelement(qn("a:effectLst"), {})
    bgPr.append(effLst)
    bg.append(bgPr)
    cSld.insert(0, bg)


def suppress_shapes(slide, names):
    """Delete shapes whose name matches a curated suppression (e.g. the stray blue bar)."""
    n = 0
    for sh in list(slide.shapes):
        if sh.name in names:
            sh._element.getparent().remove(sh._element)
            n += 1
    return n


# =====================================================================================
#  PHASE 1 — COM assembly: references copied natively, authored positions left blank
# =====================================================================================
def com_assemble(plan, entries, frag_by_slide, out_path):
    """Build the output deck via PowerPoint COM. Reference slides are InsertFromFile'd
    from their source .pptx at the correct position; authored slides get a blank slide
    placeholder (populated later by python-pptx). Curated pptx_native fragment slides:
    'replace' (s04) -> blank slide rebuilt from the fragment; 'append' (s05) -> the
    reference slide is copied, then the raster chart is stripped and native bars overlaid.
    Returns list of per-slide dicts:
    {kind:'reference'|'authored'|'fragment_replace'|'fragment_append', ...}."""
    import win32com.client as w
    slides_meta = []
    pp = w.Dispatch("PowerPoint.Application")
    try:
        out = pp.Presentations.Add(WithWindow=False)
        # PowerPoint decks start with zero slides via Add(); force TRUE widescreen
        # (13.333in x 7.5in = 960pt x 540pt). ppSlideSizeOnScreen16x9 is only 10in wide,
        # which silently rescales all authored coordinates by 1.33x.
        try:
            out.PageSetup.SlideWidth = 960
            out.PageSetup.SlideHeight = 540
        except Exception:
            pass
        pos = 0  # 0-based count of slides already in `out`
        for sl in [s for s in plan["slides"] if s.get("status") != "deleted"]:
            ref = sl.get("reference") if isinstance(sl.get("reference"), dict) else None
            if ref and ref.get("ref_id") in entries:
                e = entries[ref["ref_id"]]
                src = e["provenance"]["source_path"].replace("/", "\\")
                sidx = int(e["provenance"]["source_index"])   # 1-based
                if not os.path.exists(src):
                    raise SystemExit("EXPORT ERROR: source PPTX missing for %s: %s"
                                     % (ref["ref_id"], src))
                frag = frag_by_slide.get(sidx)
                if frag and frag.get("mode") == "replace":
                    # native rebuild from the fragment — blank slide, no source raster
                    out.Slides.Add(out.Slides.Count + 1, 12)
                    slides_meta.append({"kind": "fragment_replace", "ref_id": ref["ref_id"],
                                        "native_theme": e.get("native_theme", "light"),
                                        "source_index": sidx, "frag": frag})
                elif frag:                                 # append: copy the slide, overlay natively
                    out.Slides.InsertFromFile(src, pos, sidx, sidx)
                    slides_meta.append({"kind": "fragment_append", "ref_id": ref["ref_id"],
                                        "native_theme": e.get("native_theme", "light"),
                                        "source_index": sidx, "frag": frag})
                else:
                    out.Slides.InsertFromFile(src, pos, sidx, sidx)
                    slides_meta.append({"kind": "reference", "ref_id": ref["ref_id"],
                                        "native_theme": e.get("native_theme", "light"),
                                        "source_index": sidx})
            else:
                # blank placeholder (layout 12 = blank in the default template)
                lay = out.Slides.Count + 1
                out.Slides.Add(lay, 12)
                slides_meta.append({"kind": "authored"})
            pos += 1
        out.SaveAs(os.path.abspath(out_path))
        out.Close()
    finally:
        pp.Quit()
    return slides_meta


def _has_refs(plan):
    return any(isinstance(s.get("reference"), dict) and s["reference"].get("ref_id")
               for s in plan.get("slides", []))


def _win32_available():
    try:
        import win32com.client  # noqa: F401
        return True
    except Exception:
        return False


def pptx_assemble(plan, out_path):
    """COM-FREE assembly (macOS / any box without PowerPoint): create one blank slide per
    plan slide with python-pptx, to be populated by build_authored_slide. Authored-only
    decks export fully; verbatim reference slides (which need PowerPoint's InsertFromFile
    to copy their source) fall back to a blank and are flagged. Returns per-slide meta."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(PPTW_IN)
    prs.slide_height = Inches(PPTH_IN)
    blank = prs.slide_layouts[6]              # 6 = blank layout in the default template
    meta = []
    # Tombstoned slides (status:"deleted") stay in the plan for pin resolution but are
    # never rendered - the DOM has no section for them, so adding a slide would ship blanks.
    for sl in [s for s in plan.get("slides", []) if s.get("status") != "deleted"]:
        prs.slides.add_slide(blank)
        ref = sl.get("reference") if isinstance(sl.get("reference"), dict) else None
        meta.append({"kind": "ref_blank", "ref_id": ref.get("ref_id")} if (ref and ref.get("ref_id"))
                    else {"kind": "authored"})
    prs.save(out_path)
    return meta


def source_theme(entries):
    """Load the source presentation theme once (all demo refs share one source)."""
    from pptx import Presentation
    paths = {}
    for e in entries.values():
        p = (e.get("provenance") or {}).get("source_path")
        if not p:
            continue
        p = p.replace("/", "\\")
        if os.path.exists(p):
            paths[p] = True
    schemes = {}
    for p in paths:
        try:
            prs = Presentation(p)
            master = prs.slide_masters[0]
            theme = master.element.getroottree()  # not the theme; get theme part
            # theme part lives on the master's part rels
            th = master.part.part_related_by(
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme")
            from lxml import etree
            root = etree.fromstring(th.blob)
            schemes[p] = _theme_scheme(root)
        except Exception as ex:
            print("  [warn] theme load failed for %s: %s" % (p, ex))
            schemes[p] = {}
    return schemes


# =====================================================================================
#  PHASE 2 — python-pptx: recolor references, populate authored slides
# =====================================================================================
def export(sp, plan_path, slides_html, out_path, theme=None, do_authored=True):
    plan = load_json(plan_path)
    theme = theme or plan.get("deck", {}).get("theme", "dark")
    entries = catalog_entries(sp)
    cur = curation(sp)
    supp_by_slide = {s["slide"]: s.get("names", []) for s in cur.get("shape_suppressions", [])}
    swap_by_slide = {s["slide"]: s.get("shapes", {}) for s in cur.get("logo_swaps", [])}
    frag_by_slide = {f["slide"]: f for f in cur.get("fragments", []) if f.get("pptx_native")}
    fragdir = os.path.join(sp, "layouts", "reference-library", "fragments")

    print("EXPORT: %d slides, theme=%s" % (len([s for s in plan["slides"] if s.get("status") != "deleted"]), theme))
    t0 = time.time()
    # Use PowerPoint COM only when it's actually available AND the deck reuses verbatim
    # reference slides (which need InsertFromFile). Otherwise assemble natively with
    # python-pptx - works on macOS/Linux and for authored-only decks.
    use_com = _has_refs(plan) and _win32_available()
    if use_com:
        schemes = source_theme(entries)
        meta = com_assemble(plan, entries, frag_by_slide, out_path)
        print("  COM assembly done in %.1fs -> %s" % (time.time() - t0, out_path))
    else:
        schemes = {}
        if _has_refs(plan):
            print("  [warn] reference slides need PowerPoint (Windows) to copy verbatim; "
                  "exported as blank placeholders on this platform.")
        meta = pptx_assemble(plan, out_path)
        print("  python-pptx assembly done in %.1fs -> %s" % (time.time() - t0, out_path))

    from pptx import Presentation
    prs = Presentation(out_path)
    global SLIDE_W_EMU, SLIDE_H_EMU               # coordinate math tracks the real deck size
    SLIDE_W_EMU, SLIDE_H_EMU = prs.slide_width, prs.slide_height
    print("  deck size: %.3fin x %.3fin" % (SLIDE_W_EMU / EMU, SLIDE_H_EMU / EMU))
    manifest = None
    if do_authored:
        manifest = get_manifest(slides_html, theme)

    def _frag_html(m):
        with open(os.path.join(fragdir, m["frag"]["file"]), encoding="utf-8") as f:
            return f.read()

    rast = LogoRasterizer(sp)
    ai = 0  # authored-slide index into manifest (reference/fragment slides excluded there)
    for i, (m, slide) in enumerate(zip(meta, prs.slides)):
        kind = m["kind"]
        if kind in ("reference", "fragment_append"):
            e = entries[m["ref_id"]]
            src = e["provenance"]["source_path"].replace("/", "\\")
            sc = schemes.get(src, {})
            sidx = m["source_index"]
            names = supp_by_slide.get(sidx, [])
            if names:
                print("  s%d ref %s: suppressed %d shape(s)" % (i + 1, m["ref_id"], suppress_shapes(slide, names)))
            if theme != m["native_theme"]:
                n = recolor_slide(slide, sc, theme)
                set_background(slide, theme)
                print("  s%d ref %s: recoloured %d colour node(s) -> %s" % (i + 1, m["ref_id"], n, theme))
            else:
                print("  s%d ref %s: native theme match, copied verbatim" % (i + 1, m["ref_id"]))
            swaps = swap_by_slide.get(sidx, {})
            if swaps:
                ns = apply_logo_swaps(slide, swaps, theme, rast)
                print("  s%d ref %s: swapped %d brand logo(s) -> crisp SVG" % (i + 1, m["ref_id"], ns))
            if kind == "fragment_append":
                removed = strip_media(slide)
                man = fragment_manifest(_frag_html(m), theme, bg_hex="")   # overlay, keep copied bg
                build_authored_slide(prs, slide, man, theme)
                print("  s%d ref %s: native chart overlay (stripped %d raster media)"
                      % (i + 1, m["ref_id"], removed))
        elif kind == "fragment_replace":
            man = fragment_manifest(_frag_html(m), theme)
            build_authored_slide(prs, slide, man, theme)
            print("  s%d ref %s: native fragment rebuild (%s)" % (i + 1, m["ref_id"], m["frag"]["file"]))
        elif kind == "ref_blank":
            print("  s%d ref %s: blank placeholder (no PowerPoint to copy the source slide)" % (i + 1, m.get("ref_id")))
            # does NOT consume a manifest entry (reference slides aren't in the authored manifest)
        else:  # authored
            if do_authored and manifest and ai < len(manifest):
                build_authored_slide(prs, slide, manifest[ai], theme)
            ai += 1
    rast.close()
    prs.save(out_path)
    print("EXPORT complete in %.1fs -> %s" % (time.time() - t0, out_path))
    return out_path


# =====================================================================================
#  MANIFEST — pull the authored-slide DOM manifest from the built HTML via Playwright
# =====================================================================================
def get_manifest(slides_html, theme):
    """Drive headless Chromium, call window.__exportManifest(), return the authored-slide
    manifests (reference slides are excluded there).

    ``theme`` ('dark' | 'light') is FORCED onto the DOM (``data-theme`` on <html>) BEFORE the
    manifest is pulled, so every computed style — card fills, text colour, backgrounds — comes
    from ONE theme. Without this the export captured whatever the HTML happened to render and
    could mix a light-mode card fill with dark-mode text. Theme is REQUIRED (no default) so an
    export can never silently guess; the caller must have asked the user which one."""
    if theme not in ("dark", "light"):
        raise SystemExit("EXPORT ERROR: get_manifest(theme=%r) — theme must be 'dark' or 'light'" % theme)
    from playwright.sync_api import sync_playwright
    html = os.path.abspath(slides_html).replace("\\", "/")
    with sync_playwright() as p:
        b = p.chromium.launch()
        pg = b.new_page(viewport={"width": 1280, "height": 720}, device_scale_factor=1)
        pg.goto("file:///" + html)
        # Force the chosen theme, then let reveal/layout settle before capture.
        pg.evaluate("(t) => document.documentElement.setAttribute('data-theme', t)", theme)
        pg.wait_for_timeout(1200)
        data = pg.evaluate("async () => window.__exportManifest ? await window.__exportManifest() : null")
        b.close()
    if data is None:
        raise SystemExit("EXPORT ERROR: window.__exportManifest() unavailable in built HTML")
    return data


# ----------------------------------------------------------------- authored builders
# actual output-deck dimensions in EMU (set from the opened presentation so coordinate
# math never assumes a slide size that COM silently overrode).
SLIDE_W_EMU = int(PPTW_IN * EMU)
SLIDE_H_EMU = int(PPTH_IN * EMU)


def _emu_x(px):
    return int(px / LOGW * SLIDE_W_EMU)


def _emu_y(px):
    return int(px / LOGH * SLIDE_H_EMU)


def _hex(c):
    return (c or "FFFFFF").lstrip("#").upper()[:6] or "FFFFFF"


def build_authored_slide(prs, slide, man, theme):
    """Reconstruct one authored slide as native python-pptx objects from its manifest.
    Order: background -> shapes (z) -> images/svgs -> charts -> text (text-in-shape for
    small containers)."""
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    import base64, io

    # background
    if man.get("bg"):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(_hex(man["bg"]))

    # shapes + images/svgs, drawn in ONE z-ordered pass so a scrim that sits above the photo in
    # the HTML also sits above it here (was: all shapes, then all images -> scrims buried).
    def _draw_shape(s):
        b = s["box"]
        shp_kind = s.get("shape", "rect")
        auto = {"rect": MSO_SHAPE.RECTANGLE, "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
                "ellipse": MSO_SHAPE.OVAL}.get(shp_kind, MSO_SHAPE.RECTANGLE)
        sh = slide.shapes.add_shape(auto, _emu_x(b["x"]), _emu_y(b["y"]),
                                    _emu_x(b["w"]), _emu_y(b["h"]))
        if shp_kind == "roundRect":
            # Match the HTML's ABSOLUTE corner radius (6px brand rule). A ROUNDED_RECTANGLE's
            # adjustment is a FRACTION of the shorter side, and PowerPoint defaults it to ~0.167
            # (a pill on tall tiles). Set it from rectRadius/shorter-side so 6px stays 6px.
            shorter = max(1.0, min(b["w"], b["h"]))
            try:
                sh.adjustments[0] = max(0.0, min(0.5, s.get("rectRadius", 6) / shorter))
            except Exception:
                pass
        if s.get("gradient"):
            _apply_gradient(sh, s["gradient"])
        elif s.get("fill"):
            sh.fill.solid()
            sh.fill.fore_color.rgb = RGBColor.from_string(_hex(s["fill"]))
            if s.get("fillAlpha", 1) < 1:
                _set_alpha(sh.fill.fore_color, s["fillAlpha"])
        else:
            sh.fill.background()
        if s.get("line"):
            sh.line.color.rgb = RGBColor.from_string(_hex(s["line"]))
            sh.line.width = Pt(max(0.5, s.get("lineW", 1) * 0.75))
        else:
            sh.line.fill.background()
        sh.shadow.inherit = False
        # text-in-shape for small containers with a bound text run
        if s.get("text") and b["w"] <= 72 and b["h"] <= 72:
            tf = sh.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            _apply_runs(tf, s["text"], PP_ALIGN, Pt, RGBColor, center=True)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _draw_image(im):
        b = im["box"]
        data = im.get("data")
        if not data or not data.startswith("data:"):
            return
        try:
            raw = base64.b64decode(data.split(",", 1)[1])
            slide.shapes.add_picture(io.BytesIO(raw), _emu_x(b["x"]), _emu_y(b["y"]),
                                     _emu_x(b["w"]), _emu_y(b["h"]))
        except Exception:
            pass

    drawables = ([("shape", s.get("z", 0), s) for s in man.get("shapes", [])]
                 + [("image", im.get("z", 1 << 30), im) for im in man.get("images", [])])
    drawables.sort(key=lambda d: d[1])
    for kind, _z, obj in drawables:
        (_draw_shape if kind == "shape" else _draw_image)(obj)

    # charts
    for ch in man.get("charts", []):
        try:
            _add_chart(slide, ch)
        except Exception as ex:
            print("    [warn] chart failed: %s" % ex)

    # text boxes (skip those already emitted as text-in-shape)
    for t in man.get("texts", []):
        if t.get("inShape"):
            continue
        b = t["box"]
        tb = slide.shapes.add_textbox(_emu_x(b["x"]), _emu_y(b["y"]),
                                      _emu_x(b["w"]), _emu_y(b["h"]))
        tf = tb.text_frame
        # PowerPoint renders Montserrat a hair wider than Chrome, so any box left to re-wrap adds
        # a line and overflows the next block. Two cases turn wrapping OFF: (a) a box only one line
        # tall (e.g. "90%"), and (b) a box whose runs already carry the browser's exact soft-wrap
        # breaks (breakLine) — we reproduce those verbatim as separate paragraphs, so PowerPoint
        # must NOT re-flow them. Text boxes don't clip, so unwrapped overflow is invisible.
        lh = t.get("lineHeight")
        single_line = bool(lh) and b["h"] <= lh * 1.4
        has_breaks = any(r.get("breakLine") for r in (t.get("runs") or []))
        no_wrap = single_line or has_breaks
        tf.word_wrap = not no_wrap
        tf.auto_size = MSO_AUTO_SIZE.NONE if no_wrap else MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        from pptx.enum.text import PP_ALIGN as PA
        _apply_runs(tf, t, PA, Pt, RGBColor)

    # page-number chrome
    pn = man.get("pageNum")
    if pn:
        b = pn["box"]
        tb = slide.shapes.add_textbox(_emu_x(b["x"]), _emu_y(b["y"]),
                                      _emu_x(b["w"]) + Emu(int(0.35 * EMU)), _emu_y(b["h"]))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        r = tf.paragraphs[0].add_run()
        r.text = pn["text"]
        r.font.name = "Montserrat"
        r.font.size = Pt(pn.get("size", 14) * 0.75)
        r.font.color.rgb = RGBColor.from_string(_hex(pn.get("color", "FFFFFF")))


def _set_alpha(color, alpha):
    """Add an <a:alpha> to a solid fill colour (python-pptx has no direct setter)."""
    from pptx.oxml.ns import qn
    srgb = color._xFill.find(qn("a:srgbClr"))
    if srgb is not None:
        a = srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100000))})
        srgb.append(a)


def _apply_runs(tf, spec, PP_ALIGN, Pt, RGBColor, center=False):
    """Fill a text_frame from a manifest text spec {runs:[{text,color,bold,size,breakLine}], align,
    lineHeight}. A run flagged breakLine starts a NEW PARAGRAPH after it — that is how the browser's
    exact soft-wrap (and <br>) breaks are reproduced so PowerPoint never re-flows the text."""
    al = "center" if center else spec.get("align", "left")
    align_val = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(al, PP_ALIGN.LEFT)
    lh = spec.get("lineHeight")

    def _setup(p):
        p.alignment = align_val
        if lh:
            # Pin line spacing to the HTML's computed line-height (px->pt). PowerPoint's default
            # spacing is looser, so a multi-line block overflowed its box downward and swallowed
            # the gap to the next box (headline crowding the body). Matching the browser keeps it.
            try:
                p.line_spacing = Pt(lh * 0.75)
            except Exception:
                pass

    para = tf.paragraphs[0]
    _setup(para)
    runs = spec.get("runs") or [{"text": spec.get("text", "")}]
    for i, rspec in enumerate(runs):
        r = para.add_run()
        r.text = rspec.get("text", "")
        r.font.name = "Montserrat"
        if rspec.get("size"):
            # NOTE: manifest run `size` is ALREADY in points — pptRuns() in deck.js does the
            # px->pt conversion (fontSize * 0.75) at capture time. Do NOT multiply again here;
            # a second * 0.75 shrank every authored text run to 75% of intended (the "font
            # clash" the reveal fix appeared to cause but never did — shapes/tiles were fine).
            r.font.size = Pt(rspec["size"])
        if rspec.get("bold"):
            r.font.bold = True
        if rspec.get("italic"):
            r.font.italic = True
        if rspec.get("underline"):
            r.font.underline = True
        r.font.color.rgb = RGBColor.from_string(_hex(rspec.get("color", "FFFFFF")))
        if rspec.get("breakLine") and i < len(runs) - 1:
            para = tf.add_paragraph()
            _setup(para)


def _apply_gradient(sh, grad):
    """Give a shape a native linear <a:gradFill> from a manifest gradient {angle(css deg),
    stops:[{color,alpha,pos}]}. Used for legibility scrims (dark->clear) that a flat fill can't
    represent — a flat scrim darkens the whole slide and buries the artwork under it."""
    from pptx.oxml.ns import qn
    try:
        spPr = sh._element.spPr
        for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
            e = spPr.find(qn(tag))
            if e is not None:
                spPr.remove(e)
        gf = spPr.makeelement(qn("a:gradFill"), {})
        gs_lst = gf.makeelement(qn("a:gsLst"), {})
        stops = sorted(grad.get("stops", []), key=lambda s: s.get("pos", 0))
        for st in stops:
            pos = max(0.0, min(1.0, st.get("pos", 0)))
            gs = gs_lst.makeelement(qn("a:gs"), {"pos": str(int(pos * 100000))})
            srgb = gs.makeelement(qn("a:srgbClr"), {"val": _hex(st.get("color", "000000"))})
            srgb.append(srgb.makeelement(qn("a:alpha"),
                        {"val": str(int(max(0.0, min(1.0, st.get("alpha", 1))) * 100000))}))
            gs.append(srgb)
            gs_lst.append(gs)
        gf.append(gs_lst)
        # CSS angle (0=up, clockwise) -> OOXML angle (0=->right, clockwise), 60000ths of a degree.
        ang = (grad.get("angle", 180) - 90) % 360
        gf.append(gf.makeelement(qn("a:lin"), {"ang": str(int(ang * 60000)), "scaled": "1"}))
        # Fill must follow the geometry element in spPr's schema order — insert it right there
        # rather than appending (an append after <a:ln> would make PowerPoint repair the file).
        geom = spPr.find(qn("a:prstGeom"))
        if geom is None:
            geom = spPr.find(qn("a:custGeom"))
        if geom is not None:
            geom.addnext(gf)
        else:
            spPr.append(gf)
    except Exception:
        # fall back to whatever solid fill was already set by the caller
        pass


def _add_chart(slide, ch):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Emu
    spec = ch["spec"]
    b = ch["box"]
    cd = CategoryChartData()
    cd.categories = spec.get("labels", [])
    cd.add_series(spec.get("name", "Series 1"), spec.get("values", []))
    ct = {"doughnut": XL_CHART_TYPE.DOUGHNUT, "pie": XL_CHART_TYPE.PIE,
          "bar": XL_CHART_TYPE.COLUMN_CLUSTERED, "col": XL_CHART_TYPE.COLUMN_CLUSTERED
          }.get(spec.get("type"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    gf = slide.shapes.add_chart(ct, _emu_x(b["x"]), _emu_y(b["y"]),
                                _emu_x(b["w"]), _emu_y(b["h"]), cd)
    gf.chart.has_legend = False
    try:
        gf.chart.has_title = False
    except Exception:
        pass


# ----------------------------------------------------------------- CLI
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--skill-path", default=os.path.dirname(HERE))
    ap.add_argument("--plan", required=True)
    ap.add_argument("--slides-html", help="built presentation.html (for authored-slide manifest)")
    ap.add_argument("--out", required=True)
    ap.add_argument("--theme", required=True, choices=["dark", "light"],
                    help="REQUIRED — the single theme to export in. Always ask the user which "
                         "one before exporting; never assume the deck's authoring theme.")
    ap.add_argument("--no-authored", action="store_true", help="skip authored slides (references only)")
    a = ap.parse_args()
    export(a.skill_path, a.plan, a.slides_html, a.out, a.theme,
           do_authored=not a.no_authored)


if __name__ == "__main__":
    main()
