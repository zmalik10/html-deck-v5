#!/usr/bin/env python3
"""verify_export_macos.py — render-verify a PPTX export on macOS (no PowerPoint COM).

PowerPoint's AppleScript surface is sandbox-blocked for scripted open/save (error
-9074), so this drives KEYNOTE, which imports PPTX and exports slide images
reliably. Keynote is a different renderer than PowerPoint, but with the deck's
fonts installed it uses the same font metrics and catches the failure classes
that matter: missing/blank slides, elements that bleed past their containers,
decor that exports as stray marks, and wrapped text that merged into long lines.

Checks (exit 1 on any failure):
  1. python-pptx structural pass — slide count matches the plan's LIVE slides,
     zero empty slides.
  2. Keynote renders one PNG per slide (count must match).
  3. Renders land in --out-dir for eyeballing / diffing against the HTML.

Usage:
    python engine/verify_export_macos.py --plan <deck>/plan.json \
        --pptx "<deck>/out/deck.pptx" --out-dir <dir>

Fonts: the PPTX references Montserrat by name. Install assets/fonts/ttf/*.ttf
to ~/Library/Fonts first or every render check is measuring a substitute font.
"""
import argparse, json, os, subprocess, sys, tempfile


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--plan", required=True)
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--out-dir", default=None)
    args = ap.parse_args()

    plan = json.load(open(args.plan, encoding="utf-8"))
    live = len([s for s in plan.get("slides", []) if s.get("status") != "deleted"])
    pptx = os.path.abspath(args.pptx)
    out_dir = os.path.abspath(args.out_dir or tempfile.mkdtemp(prefix="pptx-verify-"))
    os.makedirs(out_dir, exist_ok=True)

    fails = []

    # 1) structural pass
    try:
        from pptx import Presentation
        prs = Presentation(pptx)
        n = len(prs.slides)
        empty = [i for i, s in enumerate(prs.slides, 1) if len(s.shapes) == 0]
        if n != live:
            fails.append("slide count %d != %d live plan slides (tombstones exported?)" % (n, live))
        if empty:
            fails.append("empty slides at positions %s" % empty)
        print("structural: %d slides, %d live in plan, empty: %s" % (n, live, empty or "none"))
    except Exception as e:
        fails.append("python-pptx could not open the deck: %s" % e)

    # 2) Keynote render
    script = '''
with timeout of 300 seconds
    tell application "Keynote"
        open POSIX file "%s"
        delay 5
        set d to front document
        export d as slide images to POSIX file "%s" with properties {image format:PNG, skipped slides:false}
        close d saving no
    end tell
end timeout''' % (pptx, out_dir)
    r = subprocess.run(["osascript", "-e", script], capture_output=True, text=True)
    if r.returncode != 0:
        fails.append("Keynote render failed: %s" % (r.stderr.strip() or r.stdout.strip()))
    pngs = sorted(f for f in os.listdir(out_dir) if f.lower().endswith(".png"))
    print("keynote: %d slide renders -> %s" % (len(pngs), out_dir))
    if pngs and len(pngs) != live:
        fails.append("Keynote rendered %d slides, plan has %d live" % (len(pngs), live))

    if fails:
        print("\nVERIFY FAILED:")
        for f in fails:
            print("  FAIL " + f)
        sys.exit(1)
    print("\nVERIFY OK — now EYEBALL the renders in %s against out/review.html "
          "(text overflow, bleeding images, stray decor)." % out_dir)


if __name__ == "__main__":
    main()
